"""
Template Style Engine: analyzes template pages and produces layout blueprints
for content-to-template mapping. Zero hardcoding.

Core principle: Instead of trying to "rebuild" template visuals (which python-pptx
cannot do for gradients/shadows/images), we produce a precise "blueprint" that
guides the TemplateCloner to map generated content onto the original template's
text shapes while preserving all decorations.
"""

import logging
import re
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)


def _parse_modules(markdown: str) -> List[dict]:
    """Extract content modules from markdown."""
    modules = []
    if not markdown:
        return modules
    pattern = r'###\s*([^\n]+)(.*?)(?=###|\Z)'
    for title, content in re.findall(pattern, markdown, re.DOTALL):
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        bullets = []
        for l in lines:
            for marker in ('- ', '* ', '\u2022 ', '\u2013 ', '\u2014 '):
                if l.startswith(marker):
                    bullet_text = l[len(marker):]
                    # Strip any nested bullet markers (e.g., LLM generates "- • text")
                    for inner in ('\u2022 ', '\u00b7 ', '- ', '* '):
                        if bullet_text.startswith(inner):
                            bullet_text = bullet_text[len(inner):]
                    bullets.append(bullet_text)
                    break
        modules.append({'title': title.strip(), 'bullets': bullets})
    return modules


def _shape_text(shape: dict) -> str:
    """Concatenate all text in a shape."""
    text = ""
    for p in shape.get("paragraphs", []):
        text += p.get("text", "")
    return text.strip()


def _extract_font_info(shape: dict) -> dict:
    """Extract best available font info from a shape (run-level first, then paragraph-level)."""
    default = {
        "font_name": None,
        "font_size": None,
        "bold": None,
        "italic": None,
        "underline": None,
        "color": None,
    }
    paragraphs = shape.get("paragraphs", [])
    if not paragraphs:
        return default

    # Try run-level first (most precise)
    for para in paragraphs:
        for run in para.get("runs", []):
            if any(run.get(k) for k in ["font_name", "font_size", "color"]):
                return {
                    "font_name": run.get("font_name"),
                    "font_size": run.get("font_size"),
                    "bold": run.get("bold"),
                    "italic": run.get("italic"),
                    "underline": run.get("underline"),
                    "color": run.get("color"),
                }

    # Fallback to paragraph-level font (important for placeholders)
    para = paragraphs[0]
    pf = para.get("font", {})
    if any(pf.get(k) for k in ["font_name", "font_size", "color"]):
        return {
            "font_name": pf.get("font_name"),
            "font_size": pf.get("font_size"),
            "bold": pf.get("bold"),
            "italic": pf.get("italic"),
            "underline": pf.get("underline"),
            "color": pf.get("color"),
        }

    return default


def _estimate_text_capacity(width_emu: int, height_emu: int, font_size_pt: float, is_chinese: bool = True) -> dict:
    """Estimate how much text fits in a box.
    
    Applies safety margins to account for:
    - Paragraph margins/bullet indentation not in raw width
    - Mixed Chinese/English text having wider average character width
    - PowerPoint rendering variance
    """
    pt_per_emu = 1 / 12700
    w_pt = width_emu * pt_per_emu
    h_pt = height_emu * pt_per_emu
    char_w_pt = font_size_pt * (0.6 if is_chinese else 0.5)
    # Safety margin: 15% width reduction for margins and rendering variance
    chars_per_line = max(1, int(w_pt * 0.85 / char_w_pt))
    line_h_pt = font_size_pt * 1.35
    max_lines = max(1, int(h_pt / line_h_pt))
    return {
        "chars_per_line": chars_per_line,
        "max_lines": max_lines,
        "total_chars": chars_per_line * max_lines,
    }


def _hex_from_desc(desc: str) -> Optional[str]:
    """Try to extract a hex color from text description."""
    if not desc:
        return None
    match = re.search(r'#?[0-9A-Fa-f]{6}', desc)
    if match:
        h = match.group(0)
        return h if h.startswith("#") else f"#{h}"
    mapping = {
        "navy": "#000080", "deep blue": "#00008B", "royal blue": "#4169E1",
        "tech blue": "#0078D4", "sky blue": "#87CEEB", "light blue": "#E8F4FC",
        "cyan": "#00FFFF", "white": "#FFFFFF", "black": "#000000",
        "red": "#FF0000", "green": "#008000", "orange": "#FFA500",
        "purple": "#800080", "gray": "#808080", "slate": "#708090",
        "charcoal": "#36454F",
    }
    desc_lower = desc.lower()
    for name, hex_val in sorted(mapping.items(), key=lambda x: -len(x[0])):
        if name in desc_lower:
            return hex_val
    return None


def _group_by_position(items: List[dict], key: str, tolerance: int) -> List[List[dict]]:
    sorted_items = sorted(items, key=lambda x: x.get(key, 0))
    groups = []
    for item in sorted_items:
        matched = False
        for g in groups:
            avg = sum(x.get(key, 0) for x in g) / len(g)
            if abs(item.get(key, 0) - avg) <= tolerance:
                g.append(item)
                matched = True
                break
        if not matched:
            groups.append([item])
    return groups


class TemplateLayoutMapper:
    """Analyzes a template page and produces a content-mapping blueprint."""

    def __init__(self, template_page: dict, pptx_path: str = None):
        self.layout_json = template_page.get("layout_json", {}) if isinstance(template_page, dict) else {}
        self.visual_json = template_page.get("visual_json", {}) if isinstance(template_page, dict) else {}
        self.slide_w = self.layout_json.get("slide_width", 9144000)
        self.slide_h = self.layout_json.get("slide_height", 6858000)
        self.shapes = self.layout_json.get("shapes", [])
        # Supplement table info from actual PPTX if path provided
        if pptx_path:
            self._supplement_table_info(pptx_path)
        self.blueprint = self._build_blueprint()

    def _supplement_table_info(self, pptx_path: str):
        """Read actual PPTX file to supplement table info not in layout_json.
        Extracts per-column capacity metrics (width, font size, max text length,
        max lines) to guide LLM generation and prevent overflow."""
        try:
            from pptx import Presentation
            from pptx.oxml.ns import qn
            prs = Presentation(pptx_path)
            slide_idx = self.layout_json.get("slide_index", 0)
            if slide_idx >= len(prs.slides):
                return
            slide = prs.slides[slide_idx]
            for shp in slide.shapes:
                if not shp.has_table:
                    continue
                table = shp.table
                rows = list(table.rows)
                cols = list(table.columns)
                num_cols = len(cols)
                num_rows = len(rows)
                
                # Extract headers
                headers = []
                if rows:
                    headers = [cell.text.strip() for cell in rows[0].cells]
                
                # Compute per-column metrics
                columns_info = []
                for c_idx in range(num_cols):
                    col_width = cols[c_idx].width if c_idx < len(cols) else 0
                    
                    # Gather all cells in this column
                    max_text_len = 0
                    max_lines = 1
                    font_size = None
                    for r_idx in range(num_rows):
                        if r_idx >= len(rows):
                            break
                        cells = list(rows[r_idx].cells)
                        if c_idx >= len(cells):
                            break
                        cell = cells[c_idx]
                        text = cell.text.strip()
                        cell_lines = text.split('\n')
                        max_lines = max(max_lines, len(cell_lines))
                        for line in cell_lines:
                            max_text_len = max(max_text_len, len(line))
                        
                        # Capture font size from first run (prefer data rows over header)
                        if cell.text_frame.paragraphs:
                            p = cell.text_frame.paragraphs[0]
                            if p.runs and p.runs[0].font.size:
                                font_size = p.runs[0].font.size.pt
                    
                    columns_info.append({
                        "index": c_idx,
                        "width": col_width,
                        "header": headers[c_idx] if c_idx < len(headers) else f"列{c_idx+1}",
                        "font_size": font_size,
                        "max_text_len": max_text_len,
                        "max_lines": max_lines,
                    })
                
                table_info = {
                    "rows": num_rows,
                    "cols": num_cols,
                    "headers": headers,
                    "columns": columns_info,
                }
                
                # Check if this table is already in shapes list
                found = False
                for s in self.shapes:
                    if s.get("shape_id") == shp.shape_id:
                        found = True
                        existing = s.get("table", {})
                        existing.update(table_info)
                        s["table"] = existing
                        break
                
                if not found:
                    sample_rows = []
                    for r_idx in range(1, min(num_rows, 4)):
                        if r_idx < len(rows):
                            sample_rows.append([cell.text.strip() for cell in rows[r_idx].cells])
                    
                    self.shapes.append({
                        "shape_id": shp.shape_id,
                        "name": shp.name,
                        "shape_type": str(shp.shape_type),
                        "left": shp.left,
                        "top": shp.top,
                        "width": shp.width,
                        "height": shp.height,
                        "table": {
                            **table_info,
                            "sample_rows": sample_rows,
                        },
                    })
        except Exception as e:
            logger.warning(f"Failed to supplement table info from {pptx_path}: {e}")

    def _build_blueprint(self) -> dict:
        """Analyze all shapes and build a mapping blueprint."""
        text_slots = []
        decorations = []
        tables = []
        for s in self.shapes:
            # Detect table shapes
            if s.get("table"):
                tables.append(s)
                continue
            txt = _shape_text(s)
            if txt:
                font_info = _extract_font_info(s)
                text_slots.append({
                    **s,
                    "_text": txt,
                    "_text_len": len(txt),
                    "_font": font_info,
                    "_area": s.get("width", 0) * s.get("height", 0),
                })
            else:
                decorations.append(s)

        # Classify text slots by position and text characteristics
        slots_by_type = self._classify_slots(text_slots)
        slots_by_type["tables"] = tables

        # Build content regions (groups of related shapes)
        regions = self._build_regions(slots_by_type, decorations)

        # Compute capacity for each content slot
        contents = slots_by_type.get("content", [])
        for slot in contents:
            font_size = slot["_font"].get("font_size") or 14
            slot["_capacity"] = _estimate_text_capacity(
                slot.get("width", 0), slot.get("height", 0), font_size
            )

        # Compute capacity for each label slot as well
        labels = slots_by_type.get("labels", [])
        for slot in labels:
            font_size = slot["_font"].get("font_size") or 14
            slot["_capacity"] = _estimate_text_capacity(
                slot.get("width", 0), slot.get("height", 0), font_size
            )

        # Group content slots into visual rows for contextual generation
        slots_by_type["_rows"] = self._group_slots_into_rows(contents)
        slots_by_type["_row_order"] = self._flatten_rows(slots_by_type["_rows"])

        # Extract colors
        colors = self._extract_colors(decorations, text_slots)

        # Extract slide background from visual analysis
        vision = (self.visual_json or {}).get("vision_analysis", {})
        bg_color = colors.get("background")
        if not bg_color:
            vcs = vision.get("color_scheme", {})
            bg_color = _hex_from_desc(vcs.get("background", "")) or "#FFFFFF"

        return {
            "slide_width": self.slide_w,
            "slide_height": self.slide_h,
            "slots": slots_by_type,
            "regions": regions,
            "colors": colors,
            "background_color": bg_color,
            "content_slots_count": len(slots_by_type.get("content", [])),
            "title_slot": slots_by_type.get("title"),
            "subtitle_slot": slots_by_type.get("subtitle"),
            "visual_summary": (self.visual_json or {}).get("visual_summary", ""),
            "layout_pattern": (self.visual_json or {}).get("layout_pattern", ""),
            "vision_color_scheme": vision.get("color_scheme", {}),
            "vision_shadow": vision.get("shadow_effects", ""),
            "vision_gradient": vision.get("gradient_description", ""),
            "vision_border": vision.get("border_styles", ""),
            "vision_spacing": vision.get("spacing_alignment", ""),
        }

    def _classify_slots(self, text_slots: List[dict]) -> dict:
        """Classify text slots into title, subtitle, label, content."""
        title = None
        subtitle = None
        labels = []
        contents = []

        # Sort by vertical position for stable classification
        sorted_slots = sorted(text_slots, key=lambda s: (s.get("top", 0), s.get("left", 0)))

        for slot in sorted_slots:
            text = slot["_text"]
            top = slot.get("top", 0)
            width = slot.get("width", 0)
            height = slot.get("height", 0)
            area = slot.get("_area", 0)
            name = (slot.get("name") or "").lower()

            # Skip slide number placeholders — they should never be treated as labels or content
            is_slide_number = (
                "灯片编号" in name
                or "slide number" in name
                or "slidenum" in name
            )
            if is_slide_number:
                continue

            # Label: small, short numeric text (e.g. ①②③, 1. 2. 3.)
            is_label = (
                width < self.slide_w * 0.15
                and height < self.slide_h * 0.15
                and len(text) <= 8
                and (text.isdigit() or re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d\.]+$', text))
            )

            # Title: near top, very wide, relatively short, AND short text
            is_title = (
                not is_label
                and top < self.slide_h * 0.15
                and width > self.slide_w * 0.6
                and height < self.slide_h * 0.15
                and len(text) < 80  # Titles are usually concise
            )

            # Subtitle: near top, wide, below title area, moderate text length
            is_subtitle = (
                not is_label
                and not is_title
                and top < self.slide_h * 0.28
                and width > self.slide_w * 0.5
                and len(text) < 300
            )

            # Small caption/label: single-line, low-capacity text box that is
            # clearly a tag/caption rather than a main content area.
            font_size = _extract_font_info(slot).get("font_size") or 14
            capacity = _estimate_text_capacity(width, height, font_size)
            is_small_caption = (
                not is_label
                and not is_title
                and not is_subtitle
                and capacity["total_chars"] <= 55
                and capacity["max_lines"] == 1
                and width < self.slide_w * 0.45
            )

            if is_label or is_small_caption:
                labels.append(slot)
            elif is_title and title is None:
                title = slot
            elif is_subtitle and subtitle is None and title is not None:
                subtitle = slot
            else:
                contents.append(slot)

        # If no title found, pick the topmost wide slot
        if title is None and sorted_slots:
            for slot in sorted_slots:
                if slot.get("width", 0) > self.slide_w * 0.4:
                    title = slot
                    if slot in contents:
                        contents.remove(slot)
                    break

        # Sort contents by reading order.
        # Detect if slots are mostly in a single row (similar tops) — then sort by left.
        if len(contents) >= 2:
            tops = [s.get("top", 0) for s in contents]
            heights = [s.get("height", 0) for s in contents]
            top_range = max(tops) - min(tops)
            avg_height = sum(heights) / len(heights) if heights else 0
            # If vertical spread is small relative to shape height, treat as horizontal row(s)
            if top_range < avg_height * 1.5:
                contents.sort(key=lambda s: s.get("left", 0))
            else:
                contents.sort(key=lambda s: (s.get("top", 0), s.get("left", 0)))
        else:
            contents.sort(key=lambda s: (s.get("top", 0), s.get("left", 0)))
        labels.sort(key=lambda s: (s.get("top", 0), s.get("left", 0)))

        # Filter out decorative tiny shapes (capacity <= 2 and text <= 2 chars)
        # These are usually visual ornaments (e.g. single char in a diagram node)
        # and should keep their original text, not be replaced by LLM.
        filtered_labels = []
        for lbl in labels:
            cap = lbl.get("_capacity", {})
            text = lbl.get("_text", "")
            total = cap.get("total_chars", 0)
            is_decorative = (
                total <= 2
                and len(text) <= 2
                and not (text.isdigit() or re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d\.]+$', text))
            )
            if not is_decorative:
                filtered_labels.append(lbl)
        labels = filtered_labels

        return {
            "title": title,
            "subtitle": subtitle,
            "labels": labels,
            "content": contents,
            "all": sorted_slots,
        }

    def _build_regions(self, slots_by_type: dict, decorations: List[dict]) -> List[dict]:
        """Group content shapes into spatial regions (cards, columns, rows, etc.)."""
        contents = slots_by_type.get("content", [])
        if len(contents) < 2:
            return []

        # Try to detect grid structure from content positions
        cols = _group_by_position(contents, "left", int(self.slide_w * 0.08))
        rows = _group_by_position(contents, "top", int(self.slide_h * 0.08))

        regions = []
        if len(cols) >= 2 and len(rows) >= 2:
            # Grid layout
            for i, slot in enumerate(contents):
                col = i % len(cols)
                row = i // len(cols)
                regions.append({
                    "region_id": f"cell_{row}_{col}",
                    "type": "grid_cell",
                    "slot": slot,
                    "row": row,
                    "col": col,
                })
        elif len(cols) >= 2:
            # Horizontal columns
            for i, slot in enumerate(contents):
                regions.append({
                    "region_id": f"col_{i}",
                    "type": "column",
                    "slot": slot,
                    "col": i,
                })
        elif len(rows) >= 2:
            # Vertical rows
            for i, slot in enumerate(contents):
                regions.append({
                    "region_id": f"row_{i}",
                    "type": "row",
                    "slot": slot,
                    "row": i,
                })

        return regions

    def _group_slots_into_rows(self, slots: List[dict]) -> List[List[dict]]:
        """Group content slots into visual rows based on vertical position.
        
        Uses tolerance of ~8% of slide height to cluster slots that are
        visually on the same horizontal row.
        """
        if not slots:
            return []
        tolerance = int(self.slide_h * 0.08)
        rows = _group_by_position(slots, "top", tolerance)
        # Sort each row by left position for reading order
        for row in rows:
            row.sort(key=lambda s: s.get("left", 0))
        # Sort rows by average top position
        rows.sort(key=lambda row: sum(s.get("top", 0) for s in row) / len(row))
        return rows

    @staticmethod
    def _flatten_rows(rows: List[List[dict]]) -> List[dict]:
        """Flatten row groups into a single ordered list."""
        result = []
        for row in rows:
            result.extend(row)
        return result

    def _extract_colors(self, decorations: List[dict], text_shapes: List[dict]) -> dict:
        colors = {
            "primary": None, "secondary": None, "accent": None,
            "text": None, "background": None, "card_bg": None,
            "label_bg": None, "label_text": None,
        }

        # From decoration fills
        deco_fills = []
        for d in decorations:
            fc = d.get("fill_color")
            if fc and fc != "none" and len(fc) == 6:
                deco_fills.append((fc.upper(), d.get("width", 0) * d.get("height", 0)))

        deco_fills_filtered = [(c, a) for c, a in deco_fills if c not in ("FFFFFF", "000000", "NONE")]
        if deco_fills_filtered:
            largest = max(deco_fills_filtered, key=lambda x: x[1])
            colors["card_bg"] = f"#{largest[0]}"

        from collections import Counter
        color_counts = Counter([c for c, a in deco_fills_filtered if c != largest[0]])
        most_common = color_counts.most_common(3)
        if len(most_common) >= 1:
            colors["primary"] = f"#{most_common[0][0]}"
        if len(most_common) >= 2:
            colors["secondary"] = f"#{most_common[1][0]}"
        if len(most_common) >= 3:
            colors["accent"] = f"#{most_common[2][0]}"

        # From text shape colors
        for s in text_shapes:
            font_info = s.get("_font", {})
            c = font_info.get("color")
            if c and len(c) == 6 and c.upper() not in ("FFFFFF", "000000"):
                if colors["text"] is None:
                    colors["text"] = f"#{c}"
            if c and len(c) == 6 and c.upper() == "FFFFFF":
                if colors["label_text"] is None:
                    colors["label_text"] = f"#{c}"

        # From vision analysis
        vision = (self.visual_json or {}).get("vision_analysis", {})
        vcs = vision.get("color_scheme", {})
        if isinstance(vcs, dict):
            colors["primary"] = colors["primary"] or _hex_from_desc(vcs.get("primary", ""))
            colors["secondary"] = colors["secondary"] or _hex_from_desc(vcs.get("secondary", ""))
            colors["accent"] = colors["accent"] or _hex_from_desc(vcs.get("accent", ""))
            bg = _hex_from_desc(vcs.get("background", ""))
            if bg:
                colors["background"] = colors["background"] or bg
                colors["card_bg"] = colors["card_bg"] or bg

        # Fallbacks
        if colors["primary"] is None:
            colors["primary"] = colors["card_bg"] or "#004B8D"
        if colors["secondary"] is None:
            colors["secondary"] = colors["primary"]
        if colors["accent"] is None:
            colors["accent"] = colors["primary"]
        if colors["text"] is None:
            colors["text"] = "#1A1A1A"
        if colors["background"] is None:
            colors["background"] = "#FFFFFF"
        if colors["card_bg"] is None:
            colors["card_bg"] = colors["secondary"] or "#F2F2F2"
        if colors["label_bg"] is None:
            colors["label_bg"] = colors["primary"]
        if colors["label_text"] is None:
            colors["label_text"] = "#FFFFFF"

        return colors

    def get_style_profile(self) -> dict:
        """Return a comprehensive style profile for LLM prompts."""
        bp = self.blueprint
        contents = bp["slots"].get("content", [])
        capacities = [s.get("_capacity", {}) for s in contents]
        avg_capacity = {
            "chars_per_line": int(sum(c.get("chars_per_line", 40) for c in capacities) / len(capacities)) if capacities else 40,
            "max_lines": int(sum(c.get("max_lines", 4) for c in capacities) / len(capacities)) if capacities else 4,
            "total_chars": int(sum(c.get("total_chars", 160) for c in capacities) / len(capacities)) if capacities else 160,
        } if capacities else {"chars_per_line": 40, "max_lines": 4, "total_chars": 160}

        title_font = bp["slots"]["title"]["_font"] if bp["slots"]["title"] else {}
        body_fonts = [s["_font"] for s in contents if s.get("_font")]
        body_font = body_fonts[0] if body_fonts else {}

        return {
            "slide_width": bp["slide_width"],
            "slide_height": bp["slide_height"],
            "title_style": title_font,
            "body_style": body_font,
            "colors": bp["colors"],
            "background_color": bp["background_color"],
            "content_slots_count": bp["content_slots_count"],
            "capacity": avg_capacity,
            "visual_summary": bp["visual_summary"],
            "layout_pattern": bp["layout_pattern"],
            "vision_color_scheme": bp["vision_color_scheme"],
            "vision_shadow": bp["vision_shadow"],
            "vision_gradient": bp["vision_gradient"],
            "vision_border": bp["vision_border"],
            "vision_spacing": bp["vision_spacing"],
        }

    def get_layout_hint(self) -> str:
        """Return a human-readable layout constraint hint."""
        bp = self.blueprint
        contents = bp["slots"].get("content", [])
        title_slot = bp["slots"].get("title")
        subtitle_slot = bp["slots"].get("subtitle")

        parts = []
        if title_slot:
            tf = title_slot.get("_font", {})
            parts.append(f"Title: font={tf.get('font_name')}, size={tf.get('font_size')}pt")

        if subtitle_slot:
            parts.append("Has subtitle slot")

        if contents:
            capacities = [s.get("_capacity", {}) for s in contents]
            avg_total = int(sum(c.get("total_chars", 0) for c in capacities) / len(capacities)) if capacities else 0
            parts.append(f"Content slots: {len(contents)}, avg capacity: ~{avg_total} chars each")

        if bp.get("regions"):
            region_types = set(r["type"] for r in bp["regions"])
            parts.append(f"Layout type: {', '.join(region_types)}")

        return "; ".join(parts) if parts else "Custom layout."

    def generate_content_mapping(self, slide_data: dict) -> dict:
        """Generate a mapping of which content module goes to which template text slot."""
        bp = self.blueprint
        contents = bp["slots"].get("content", [])
        labels = bp["slots"].get("labels", [])

        mapping = {
            "slide_number": slide_data.get("slide_number", 1),
            "title": slide_data.get("title", ""),
            "subtitle": slide_data.get("subtitle", ""),
            "slots": bp["slots"],  # Include slots for TemplateCloner to find shapes
            "slot_mappings": [],
            "label_mappings": [],
            "table_mappings": slide_data.get("_table_mappings", []),
        }

        # If normalizer already computed slot mappings, reuse them
        precomputed = slide_data.get("_slot_mappings")
        if precomputed:
            mapping["slot_mappings"] = precomputed
        else:
            # Fallback: compute from markdown using row-ordered slots
            modules = _parse_modules(slide_data.get("markdown_content", ""))
            row_order = bp["slots"].get("_row_order", contents)

            for i, slot in enumerate(row_order):
                cap = slot.get("_capacity", {})
                max_lines = cap.get("max_lines", 4)
                max_chars = cap.get("chars_per_line", 40)

                if i < len(modules):
                    mod = modules[i]
                    text = mod["title"]
                    for b in mod["bullets"][:max(1, max_lines - 1)]:
                        if len(b) > max_chars:
                            b = b[:max_chars - 1] + "…"
                        text += "\n• " + b
                else:
                    text = ""

                mapping["slot_mappings"].append({
                    "shape_name": slot.get("name"),
                    "shape_id": slot.get("shape_id"),
                    "text": text,
                    "capacity": cap,
                })

        # Build label mappings from generated content or fallback to original text
        precomputed_labels = slide_data.get("labels_content", [])
        if precomputed_labels and labels:
            for i, label in enumerate(labels):
                text = precomputed_labels[i] if i < len(precomputed_labels) else ""
                mapping["label_mappings"].append({
                    "shape_name": label.get("name"),
                    "shape_id": label.get("shape_id"),
                    "text": text,
                })
        elif labels:
            # Fallback: keep original text for non-digit labels, keep digits as-is
            for i, label in enumerate(labels):
                orig = label.get("_text", "")
                if orig.isdigit() or re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d\.]+$', orig):
                    text = orig
                else:
                    text = orig  # preserve original meaningful text
                mapping["label_mappings"].append({
                    "shape_name": label.get("name"),
                    "shape_id": label.get("shape_id"),
                    "text": text,
                })

        return mapping


# Backward compatibility alias
class TemplateStyleEngine(TemplateLayoutMapper):
    """Backward-compatible alias for TemplateLayoutMapper."""
    pass

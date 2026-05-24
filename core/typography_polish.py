"""
Typography & slot-hierarchy polish (PR-Q2G) — conservative post-render layer.

Only adjusts font SIZE relationships and (for confidently-numeric table columns) paragraph
alignment. Never touches fills, borders, colors, chart objects, or shapes whose role/size is
uncertain (e.g. inherited size = None). Produces audits; reports every change.

Two passes:
  A. normalize_table_typography  — header font >= body font; center confidently-numeric cols.
  B. audit_and_fix_text_hierarchy — body text never larger than the slide title (cap only).
"""

import logging
import re

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

logger = logging.getLogger(__name__)

_NUMERIC = re.compile(r"^\s*[-+]?[\d,]+(?:\.\d+)?\s*[%％]?\s*$")


def _is_title_shape(shape) -> bool:
    name = (shape.name or "")
    return ("标题" in name) or ("title" in name.lower())


def _max_run_pt(text_frame):
    """Largest explicit run font size (pt) in a text frame, or None if all inherited."""
    best = None
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                pt = r.font.size.pt
                best = pt if best is None else max(best, pt)
    return best


def _min_run_pt(text_frame):
    best = None
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                pt = r.font.size.pt
                best = pt if best is None else min(best, pt)
    return best


def normalize_table_typography(prs, log) -> None:
    """Ensure header font >= body font; center confidently-numeric body columns.
    Conservative: only acts on explicit run sizes; preserves fills/borders/colors."""
    for sidx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not getattr(sh, "has_table", False):
                continue
            table = sh.table
            rows = list(table.rows)
            if len(rows) < 2:
                continue
            ncols = len(table.columns)
            # header font band = min explicit size across header cells
            header_sizes = [_min_run_pt(c.text_frame) for c in rows[0].cells]
            header_sizes = [s for s in header_sizes if s is not None]
            header_min = min(header_sizes) if header_sizes else None
            header_align = []
            for c in rows[0].cells:
                a = c.text_frame.paragraphs[0].alignment if c.text_frame.paragraphs else None
                header_align.append(str(a).split(".")[-1] if a is not None else "inherit")

            # A1: cap any body run larger than the header band down to header_min
            capped = 0
            if header_min is not None:
                for r in rows[1:]:
                    for c in r.cells:
                        for p in c.text_frame.paragraphs:
                            for run in p.runs:
                                if run.font.size is not None and run.font.size.pt > header_min:
                                    run.font.size = Pt(header_min)
                                    capped += 1

            # A2: center confidently-numeric body columns
            numeric_cols = []
            for ci in range(ncols):
                vals = []
                for r in rows[1:]:
                    if ci < len(r.cells):
                        vals.append(r.cells[ci].text.strip())
                vals = [v for v in vals if v]
                if vals and all(_NUMERIC.match(v) for v in vals):
                    numeric_cols.append(ci)
                    for r in rows[1:]:
                        if ci < len(r.cells):
                            for p in r.cells[ci].text_frame.paragraphs:
                                p.alignment = PP_ALIGN.CENTER

            # A3: center the header row (horizontal + vertical). Safe — alignment only.
            for c in rows[0].cells:
                for p in c.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                try:
                    c.vertical_anchor = MSO_ANCHOR.MIDDLE
                except Exception:
                    pass

            # A4: center the first column when it is a category/index/key column
            # (body cells mostly short, non-numeric labels — not a numeric data column).
            first_col_centered = False
            if ncols >= 2 and 0 not in numeric_cols:
                vals = [rows[r].cells[0].text.strip() for r in range(1, len(rows)) if rows[r].cells]
                vals = [v for v in vals if v]
                short_labels = vals and all(len(v) <= 12 for v in vals)
                if short_labels:
                    for r in rows[1:]:
                        if r.cells:
                            for p in r.cells[0].text_frame.paragraphs:
                                p.alignment = PP_ALIGN.CENTER
                            try:
                                r.cells[0].vertical_anchor = MSO_ANCHOR.MIDDLE
                            except Exception:
                                pass
                    first_col_centered = True

            log["tables"].append({
                "slide": sidx, "shape": sh.name, "rows": len(rows), "cols": ncols,
                "header_min_pt": header_min, "header_align_before": header_align,
                "header_centered": True, "first_col_centered": first_col_centered,
                "body_runs_capped": capped, "numeric_cols_centered": numeric_cols,
            })


def audit_and_fix_text_hierarchy(prs, log) -> None:
    """Audit per-slide font-size bands; cap body text whose explicit size exceeds the slide
    title size. Skips tables, charts, and role/size-uncertain shapes."""
    for sidx, slide in enumerate(prs.slides, start=1):
        title_pt = None
        bodies = []  # (shape, max_pt)
        bands = []
        for sh in slide.shapes:
            if getattr(sh, "has_table", False) or getattr(sh, "has_chart", False):
                continue
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            mx = _max_run_pt(sh.text_frame)
            role = "title" if _is_title_shape(sh) else "body"
            bands.append({"shape": sh.name, "role": role, "max_pt": mx})
            if role == "title" and mx is not None:
                title_pt = mx if title_pt is None else max(title_pt, mx)
            elif role == "body" and mx is not None:
                bodies.append((sh, mx))
        changed = 0
        if title_pt is not None:
            for sh, mx in bodies:
                if mx > title_pt:  # body larger than title → clear inversion (cap to title)
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size is not None and r.font.size.pt > title_pt:
                                r.font.size = Pt(title_pt)
                                changed += 1
        log["slides"].append({"slide": sidx, "title_pt": title_pt,
                              "bands": bands, "runs_capped": changed})
        log["changed_shapes"] += (1 if changed else 0)


def typography_polish(pptx_path: str) -> dict:
    from pptx import Presentation
    table_log = {"tables": []}
    text_log = {"slides": [], "changed_shapes": 0}
    prs = Presentation(pptx_path)
    try:
        normalize_table_typography(prs, table_log)
    except Exception as e:
        logger.warning(f"table typography failed: {e}")
    try:
        audit_and_fix_text_hierarchy(prs, text_log)
    except Exception as e:
        logger.warning(f"text hierarchy failed: {e}")
    prs.save(pptx_path)
    return {
        "table_audit": table_log["tables"],
        "table_count": len(table_log["tables"]),
        "table_body_runs_capped": sum(t.get("body_runs_capped", 0) for t in table_log["tables"]),
        "text_audit": text_log["slides"],
        "text_changed_shapes": text_log["changed_shapes"],
    }

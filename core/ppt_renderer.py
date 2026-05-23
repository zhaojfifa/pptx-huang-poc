"""
PPT Renderer: renders structured layout data into a PPTX file using python-pptx.
Supports both template-clone mode (preserves template visuals) and custom-build mode.
"""

import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from config.settings import PREVIEW_DIR
from skills.skill_mermaid.mermaid_skill import MermaidSkill

logger = logging.getLogger(__name__)


def parse_color(color_str: str):
    """Parse hex color string to RGBColor."""
    if not color_str:
        return RGBColor(0, 0, 0)
    color_str = color_str.strip("#")
    if len(color_str) == 6:
        try:
            return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
        except ValueError:
            return RGBColor(0, 0, 0)
    if len(color_str) == 3:
        try:
            return RGBColor(int(color_str[0] * 2, 16), int(color_str[1] * 2, 16), int(color_str[2] * 2, 16))
        except ValueError:
            return RGBColor(0, 0, 0)
    return RGBColor(0, 0, 0)


def emu_to_inches(emu: int) -> float:
    return emu / 914400


class PPTRenderer:
    def __init__(self):
        self.mermaid = MermaidSkill()

    def render_from_template(self, slides_data: list, template_path: str, output_path: str, blueprints: list = None) -> str:
        """Render by cloning a template and replacing text content."""
        from core.template_cloner import TemplateCloner
        cloner = TemplateCloner()
        return cloner.clone_and_fill(template_path, slides_data, output_path, blueprints=blueprints)

    def render(self, layouts: list, output_path: str) -> str:
        prs = Presentation()

        for layout in layouts:
            # Use template-derived slide dimensions if available
            slide_w_emu = layout.get("slide_width", 9144000)
            slide_h_emu = layout.get("slide_height", 6858000)
            prs.slide_width = Emu(slide_w_emu)
            prs.slide_height = Emu(slide_h_emu)

            # python-pptx requires a layout; use blank if available, else last layout
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            try:
                slide = prs.slides.add_slide(blank_layout)
            except Exception:
                slide = prs.slides.add_slide(prs.slide_layouts[0])

            # Background
            bg_color = layout.get("background_color", "#FFFFFF")
            if bg_color:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = parse_color(bg_color)

            for shape_def in layout.get("shapes", []):
                self._add_shape(slide, shape_def)

        prs.save(output_path)
        logger.info(f"PPTX saved to {output_path}")
        return output_path

    def _add_shape(self, slide, shape_def: dict):
        shape_type = shape_def.get("type")
        left = Emu(shape_def.get("left", 0))
        top = Emu(shape_def.get("top", 0))
        width = Emu(shape_def.get("width", 9144000))
        height = Emu(shape_def.get("height", 6858000))

        if shape_type == "title":
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = shape_def.get("text", "")
            p.alignment = PP_ALIGN.LEFT
            p.font.name = shape_def.get("font") or "Arial"
            p.font.size = Pt(shape_def.get("font_size", 32))
            p.font.bold = shape_def.get("bold", True)
            p.font.italic = shape_def.get("italic", False)
            p.font.underline = shape_def.get("underline", False)
            p.font.color.rgb = parse_color(shape_def.get("color", "#000000"))

        elif shape_type == "content":
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            tf.word_wrap = True
            text = shape_def.get("text", "")
            lines = text.split("\n")
            if lines:
                p = tf.paragraphs[0]
                p.text = lines[0]
                p.alignment = PP_ALIGN.LEFT
                p.font.name = shape_def.get("font") or "Arial"
                p.font.size = Pt(shape_def.get("font_size", 18))
                p.font.bold = shape_def.get("bold", False)
                p.font.italic = shape_def.get("italic", False)
                p.font.underline = shape_def.get("underline", False)
                p.font.color.rgb = parse_color(shape_def.get("color", "#333333"))
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line.lstrip("- ").lstrip("* ").lstrip("• ")
                    p.level = 0 if not line.startswith(("  ", "\t")) else 1
                    p.font.name = shape_def.get("font") or "Arial"
                    p.font.size = Pt(shape_def.get("font_size", 18))
                    p.font.bold = shape_def.get("bold", False)
                    p.font.italic = shape_def.get("italic", False)
                    p.font.underline = shape_def.get("underline", False)
                    p.font.color.rgb = parse_color(shape_def.get("color", "#333333"))

        elif shape_type == "background":
            shape = slide.shapes.add_shape(1, left, top, width, height)
            fill_color = shape_def.get("fill_color", "#FFFFFF")
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = parse_color(fill_color)
            else:
                shape.fill.background()
            # Support border/line color
            line_color = shape_def.get("line_color")
            if line_color:
                shape.line.color.rgb = parse_color(line_color)
                line_width = shape_def.get("line_width")
                if line_width:
                    shape.line.width = Pt(line_width)
            else:
                shape.line.fill.background()

        elif shape_type == "image_placeholder":
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Image: {shape_def.get('description', '')}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(shape_def.get("font_size", 14))
            p.font.color.rgb = RGBColor(150, 150, 150)
            # Add a border rect
            border = slide.shapes.add_shape(1, left, top, width, height)
            border.fill.background()
            border.line.color.rgb = RGBColor(200, 200, 200)

        elif shape_type == "chart":
            mermaid_code = shape_def.get("mermaid_code", "")
            if mermaid_code:
                try:
                    img_path = self.mermaid.render(mermaid_code)
                    if img_path.lower().endswith(('.png', '.jpg', '.jpeg', '.svg')):
                        slide.shapes.add_picture(img_path, left, top, width, height)
                    else:
                        raise RuntimeError(f"Mermaid output is not an image: {img_path}")
                except Exception as e:
                    logger.warning(f"Mermaid chart unavailable, rendering placeholder: {e}")
                    shape = slide.shapes.add_shape(1, left, top, width, height)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
                    shape.line.color.rgb = RGBColor(200, 200, 200)
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "[图表占位符]"
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(150, 150, 150)

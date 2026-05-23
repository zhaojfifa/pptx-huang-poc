"""
Template Cloner: clones template slides and fills them with generated content.
Preserves all visual decorations (backgrounds, shapes, colors, fonts, gradients, shadows).

Supports:
- Complex layouts (cards, lists, columns, freeform)
- Cycling through template pages when content slides > template slides
- Blueprint-guided precise text mapping
"""

import logging
import re
from shutil import copy as shutil_copy

from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)


def _extract_modules(markdown_content: str) -> list:
    """Extract content modules from markdown."""
    modules = []
    if not markdown_content:
        return modules
    pattern = r'###\s*([^\n]+)(.*?)(?=###|\Z)'
    for title, content in re.findall(pattern, markdown_content, re.DOTALL):
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        bullets = []
        for l in lines:
            for marker in ('- ', '* ', '\u2022 ', '\u2013 ', '\u2014 '):
                if l.startswith(marker):
                    bullet_text = l[len(marker):]
                    # Strip any nested bullet markers (e.g., LLM generates "- • text")
                    for inner in ('\u2022 ', '\u00b7 ', '- ', '* '):
                        if bullet_text.startswith(inner):
                            bullet_text = bullet_text[len(inner):]
                    bullets.append(bullet_text)
                    break
        modules.append({'title': title.strip(), 'bullets': bullets})
    return modules


def _copy_run_format(target_run, source_run):
    """Copy font formatting from source run to target run.
    Uses direct XML manipulation to avoid python-pptx creating empty
    <a:solidFill/> elements which cause color glitches in PowerPoint."""
    from pptx.oxml.ns import qn
    from copy import deepcopy

    if source_run.font.name:
        target_run.font.name = source_run.font.name
    if source_run.font.size:
        target_run.font.size = source_run.font.size
    if source_run.font.bold is not None:
        target_run.font.bold = source_run.font.bold
    if source_run.font.italic is not None:
        target_run.font.italic = source_run.font.italic
    if source_run.font.underline is not None:
        target_run.font.underline = source_run.font.underline
    try:
        src_rPr = source_run._r.find(qn('a:rPr'))
        if src_rPr is not None:
            tgt_rPr = target_run._r.get_or_add_rPr()
            # Remove existing color fills from target
            for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill', 'a:blipFill'):
                el = tgt_rPr.find(qn(tag))
                if el is not None:
                    tgt_rPr.remove(el)
            # Copy source color fill if present
            for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill', 'a:blipFill'):
                src_el = src_rPr.find(qn(tag))
                if src_el is not None:
                    tgt_rPr.append(deepcopy(src_el))
                    break
    except Exception:
        pass


def _set_shape_text(shape, text: str):
    """Replace shape text while preserving all original paragraph and run formatting.
    Never creates blank paragraphs; clones existing paragraphs to preserve pPr."""
    from pptx.oxml.ns import qn

    def _has_bullet(p):
        pPr = p._p.find(qn('a:pPr'))
        if pPr is None:
            return False
        for tag in ('a:buChar', 'a:buAutoNum', 'a:buBlip'):
            if pPr.find(qn(tag)) is not None:
                return True
        return False

    def _clean_run_rPr(run):
        """Remove empty <a:solidFill/> elements that cause color glitches."""
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is not None and len(solidFill) == 0:
                rPr.remove(solidFill)

    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = text.split('\n') if text else ['']
    paragraphs = list(tf.paragraphs)

    # Find a reference run with explicit formatting to copy for empty paragraphs
    ref_run = None
    for p in paragraphs:
        for r in p.runs:
            if r.text.strip():
                ref_run = r
                break
        if ref_run:
            break

    # If we need more paragraphs, clone the last paragraph's XML to preserve pPr
    while len(paragraphs) < len(lines):
        from copy import deepcopy
        last_p = paragraphs[-1]
        new_p_el = deepcopy(last_p._p)
        # Remove all runs from the clone
        for r_el in list(new_p_el.findall(qn('a:r'))):
            new_p_el.remove(r_el)
        tf._element.append(new_p_el)
        paragraphs = list(tf.paragraphs)

    # Set text for each paragraph, preserving existing run formatting
    for i, p in enumerate(paragraphs):
        line = lines[i] if i < len(lines) else ''
        # Strip manual bullet markers if paragraph already has bullet formatting
        if _has_bullet(p):
            # Strip up to two levels of bullet markers to handle nested bullets
            for _ in range(2):
                stripped = False
                for marker in ('• ', '- ', '* ', '· '):
                    if line.startswith(marker):
                        line = line[len(marker):]
                        stripped = True
                        break
                if not stripped:
                    break
        else:
            # Even without bullet formatting, strip leading bullet markers
            # that may have been incorrectly added by the normalizer
            for _ in range(2):
                stripped = False
                for marker in ('• ', '- ', '* ', '· '):
                    if line.startswith(marker):
                        line = line[len(marker):]
                        stripped = True
                        break
                if not stripped:
                    break
        # Fix abnormal line spacing (e.g., placeholder paragraphs with 7-8% spacing)
        # that would cause text to overlap.
        ls = p.line_spacing
        if ls is not None and (ls < 0.2 or ls > 2.0):
            p.line_spacing = 1.0

        # Pick the best run to keep: prefer the one with longest/most meaningful text
        # to preserve correct font color (avoid keeping an empty white-format run).
        best_run_idx = 0
        if p.runs:
            best_len = -1
            for idx, r in enumerate(p.runs):
                r_len = len(r.text.strip())
                if r_len > best_len:
                    best_len = r_len
                    best_run_idx = idx

        while len(p.runs) > 1:
            # Remove runs other than the best one
            rem_idx = 1 if best_run_idx == 0 else 0
            p._p.remove(p.runs[rem_idx]._r)
            if rem_idx < best_run_idx:
                best_run_idx -= 1

        if p.runs:
            p.runs[0].text = line
            # If the kept run was empty, copy format from ref_run
            if best_len == 0 and ref_run:
                _copy_run_format(p.runs[0], ref_run)
            # Clean up any empty solidFill that may have been created
            _clean_run_rPr(p.runs[0])
        else:
            # Empty paragraph: add text then copy reference formatting
            p.text = line
            if ref_run and p.runs:
                _copy_run_format(p.runs[0], ref_run)
                _clean_run_rPr(p.runs[0])

    # Remove excess paragraphs
    while len(tf.paragraphs) > len(lines):
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)


class TemplateCloner:
    def _fill_shape_text(self, shape, text: str, orig_text: str = "", capacity: dict = None):
        """设置 shape 文本；若模板原文字数>30且填充文字超过 capacity，字号减小 1pt（仅一次）。"""
        orig_len = len(orig_text.strip()) if orig_text else 0
        if orig_len == 0 and shape.has_text_frame:
            orig_len = len(shape.text_frame.text.strip())

        _set_shape_text(shape, text)

        if orig_len > 30 and capacity:
            cap_total = capacity.get("total_chars", 0)
            if len(text) > cap_total:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            try:
                                r.font.size = Pt(max(8, r.font.size.pt - 1))
                            except Exception:
                                pass

    def clone_and_fill(
        self,
        template_path: str,
        slides_data: list,
        output_path: str,
        blueprints: list = None,
    ) -> str:
        """
        Clone template and fill with generated content.

        Args:
            template_path: Path to the template .pptx file
            slides_data: List of slide content dicts
            output_path: Output path for the generated .pptx
            blueprints: Optional list of layout blueprints (from TemplateLayoutMapper)
                        to guide precise text-to-shape mapping.
        """
        shutil_copy(template_path, output_path)
        prs = Presentation(output_path)
        template_slides = list(prs.slides)
        total_template = len(template_slides)

        if total_template == 0:
            raise RuntimeError(f"Template has no slides: {template_path}")

        # Pre-duplicate all needed slides first (before any filling)
        needed_slides = []
        for i, _ in enumerate(slides_data):
            template_idx = i % total_template
            if i < total_template:
                needed_slides.append(template_slides[template_idx])
            else:
                source_slide = template_slides[template_idx]
                new_slide = self._duplicate_slide(prs, source_slide)
                needed_slides.append(new_slide)

        filled = 0
        for i, slide_data in enumerate(slides_data):
            template_idx = i % total_template
            slide = needed_slides[i]

            blueprint = None
            if blueprints and template_idx < len(blueprints):
                blueprint = blueprints[template_idx]

            self._fill_slide(slide, slide_data, blueprint)
            filled += 1

        prs.save(output_path)
        logger.info(f"Template cloned: {filled} slides filled from {template_path}")
        return output_path

    def _duplicate_slide(self, prs, source_slide):
        """Duplicate a slide in the presentation."""
        from copy import deepcopy
        slide_layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        # Remove auto-generated placeholder shapes from the new slide
        for shape in list(new_slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)
        # Copy all shapes from source to new slide
        for shape in source_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        return new_slide

    def _fill_slide(self, slide, slide_data: dict, blueprint: dict = None):
        """Fill a single slide with content."""
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        markdown = slide_data.get("markdown_content", "")
        modules = _extract_modules(markdown)

        if blueprint:
            self._fill_with_blueprint(slide, title, subtitle, modules, blueprint, slide_data)
        else:
            self._fill_with_heuristics(slide, title, subtitle, markdown, modules)

    def _fill_with_blueprint(self, slide, title: str, subtitle: str, modules: list, blueprint: dict, slide_data: dict = None):
        """Use blueprint to precisely map content to template shapes."""
        slots = blueprint.get("slots", {})

        # Fill title
        title_slot = slots.get("title")
        if title_slot and title:
            shape = self._find_shape_by_id_or_name(slide, title_slot.get("shape_id"), title_slot.get("name"))
            if shape:
                self._fill_shape_text(shape, title, title_slot.get("_text", ""), title_slot.get("_capacity"))

        # Fill subtitle
        subtitle_slot = slots.get("subtitle")
        if subtitle_slot:
            shape = self._find_shape_by_id_or_name(slide, subtitle_slot.get("shape_id"), subtitle_slot.get("name"))
            if shape:
                self._fill_shape_text(shape, subtitle, subtitle_slot.get("_text", ""), subtitle_slot.get("_capacity"))

        # Use pre-computed slot_mappings if available (from generate_content_mapping)
        slot_mappings = blueprint.get("slot_mappings", [])
        if slot_mappings:
            for slot_map in slot_mappings:
                shape = self._find_shape_by_id_or_name(
                    slide, slot_map.get("shape_id"), slot_map.get("shape_name")
                )
                if shape:
                    orig_text = shape.text_frame.text if shape.has_text_frame else ""
                    self._fill_shape_text(shape, slot_map.get("text", ""), orig_text, slot_map.get("capacity"))
        else:
            # Fallback: direct content slot mapping
            content_slots = slots.get("content", [])
            for idx, slot_info in enumerate(content_slots):
                shape = self._find_shape_by_id_or_name(slide, slot_info.get("shape_id"), slot_info.get("name"))
                if not shape:
                    continue
                if idx < len(modules):
                    mod = modules[idx]
                    cap = slot_info.get("_capacity", {})
                    max_lines = cap.get("max_lines", 4)
                    max_chars = cap.get("chars_per_line", 40)
                    text = mod["title"]
                    for b in mod["bullets"][:max(1, max_lines - 1)]:
                        if len(b) > max_chars:
                            b = b[:max_chars - 1] + "…"
                        text += "\n• " + b
                    self._fill_shape_text(shape, text, slot_info.get("_text", ""), cap)
                else:
                    _set_shape_text(shape, "")

        # Fill tables if present
        tables = slots.get("tables", [])
        if tables and slide_data:
            table_mappings = slide_data.get("_table_mappings", [])
            for tbl_slot in tables:
                shape = self._find_shape_by_id_or_name(
                    slide, tbl_slot.get("shape_id"), tbl_slot.get("name")
                )
                if not shape or not shape.has_table:
                    continue
                # Find matching table mapping
                tbl_map = None
                for tm in table_mappings:
                    if tm.get("shape_id") == tbl_slot.get("shape_id") or tm.get("shape_name") == tbl_slot.get("name"):
                        tbl_map = tm
                        break
                if tbl_map:
                    self._fill_table(shape, tbl_map)

        # Update labels
        label_mappings = blueprint.get("label_mappings", [])
        if label_mappings:
            for lbl_map in label_mappings:
                shape = self._find_shape_by_id_or_name(
                    slide, lbl_map.get("shape_id"), lbl_map.get("shape_name")
                )
                if shape:
                    text = lbl_map.get("text", "")
                    # Skip empty text to avoid clearing meaningful original labels
                    if text:
                        self._fill_shape_text(shape, text)
        else:
            labels = slots.get("labels", [])
            for idx, label_slot in enumerate(labels):
                shape = self._find_shape_by_id_or_name(slide, label_slot.get("shape_id"), label_slot.get("name"))
                if shape:
                    orig = label_slot.get("_text", "")
                    if orig.isdigit():
                        pass
                    else:
                        # Preserve original label/caption text (e.g. "湛江钢铁", "上半年分红")
                        # instead of replacing with sequential numbers.
                        self._fill_shape_text(shape, orig, orig)

    def _fill_table(self, shape, table_mapping: dict):
        """Fill a table shape with generated data, adjusting rows/cols as needed.
        Preserves original cell styling (colors, fonts, borders) and only truncates
        text that exceeds the template's proven capacity (with buffer for multi-line).
        Deletes excess rows/cols instead of leaving them blank."""
        from pptx.oxml.ns import qn
        
        table = shape.table
        headers = table_mapping.get("headers", [])
        rows_data = table_mapping.get("rows", [])
        columns_info = table_mapping.get("columns", [])
        
        target_rows = 1 + len(rows_data)  # header + data rows
        target_cols = len(headers) if headers else 1
        
        # Adjust rows: add if needed, delete excess
        current_rows = len(table.rows)
        if target_rows > current_rows:
            for _ in range(target_rows - current_rows):
                table.add_row()
        elif target_rows < current_rows:
            # Delete excess rows from the bottom (keep header + data)
            tbl_el = table._tbl
            tr_elements = tbl_el.findall(qn('a:tr'))
            # Remove from the end; keep first `target_rows` rows
            for tr in reversed(tr_elements[target_rows:]):
                tbl_el.remove(tr)
        
        # Adjust columns: add if needed, delete excess
        current_cols = len(table.columns)
        if target_cols > current_cols:
            avg_width = int(shape.width / current_cols) if current_cols > 0 else 1000000
            for _ in range(target_cols - current_cols):
                table.add_column(avg_width)
        elif target_cols < current_cols:
            # Delete excess columns from every row
            tbl_el = table._tbl
            for tr in tbl_el.findall(qn('a:tr')):
                tc_elements = tr.findall(qn('a:tc'))
                for tc in reversed(tc_elements[target_cols:]):
                    tr.remove(tc)
        
        def _set_cell_text(cell, text, preserve_font_size=True):
            """Set cell text while preserving font size."""
            orig_font_size = None
            if cell.text_frame.paragraphs:
                p = cell.text_frame.paragraphs[0]
                if p.runs:
                    orig_font_size = p.runs[0].font.size
            
            cell.text = str(text) if text is not None else ""
            
            if preserve_font_size and orig_font_size:
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = orig_font_size
        
        def _truncate_cell_text(text, col_info):
            """Only truncate if text exceeds template-proven capacity with buffer.
            Allows multi-line wrapping as long as total chars are within bounds."""
            if not col_info:
                return text
            max_total = col_info.get("max_total_chars", 100)
            if len(text) <= max_total:
                return text
            return text[:max_total - 1] + "…"
        
        # Re-fetch rows after deletion/addition
        table_rows = list(table.rows)
        
        # Fill header row
        if headers and table_rows:
            for c_idx, header in enumerate(headers):
                if c_idx < len(table_rows[0].cells):
                    cell = table_rows[0].cells[c_idx]
                    col_info = columns_info[c_idx] if c_idx < len(columns_info) else None
                    truncated = _truncate_cell_text(str(header) if header is not None else "", col_info)
                    _set_cell_text(cell, truncated)
        
        # Fill data rows
        for r_idx, row_data in enumerate(rows_data):
            table_row_idx = r_idx + 1
            if table_row_idx >= len(table_rows):
                break
            for c_idx, cell_text in enumerate(row_data):
                if c_idx >= len(table_rows[table_row_idx].cells):
                    break
                cell = table_rows[table_row_idx].cells[c_idx]
                col_info = columns_info[c_idx] if c_idx < len(columns_info) else None
                truncated = _truncate_cell_text(str(cell_text) if cell_text is not None else "", col_info)
                _set_cell_text(cell, truncated)
        
        logger.info(
            f"Filled table {shape.name}: {target_cols} cols x {len(rows_data)} data rows"
        )

    def _fill_with_heuristics(self, slide, title: str, subtitle: str, markdown: str, modules: list):
        """Fallback: use position/size heuristics to classify and fill shapes."""
        slots = {"title": [], "subtitle": [], "content": [], "label": []}
        slide_h = slide.part._element.getparent().getparent().attrib.get('cy', 6858000)
        # slide height in EMU from presentation props if available
        try:
            slide_h = slide.part._element.getparent().getparent().attrib.get('cy', 6858000)
            slide_h = int(slide_h)
        except Exception:
            slide_h = 6858000

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            text = full_text.strip()
            if not text:
                continue

            area = shape.width * shape.height
            name = (shape.name or "").lower()

            is_label = (
                shape.width < 700_000
                and shape.height < 700_000
                and len(text) <= 6
                and (text.isdigit() or re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d\.]+$', text))
            )
            is_title = (
                shape.top < slide_h * 0.12
                and shape.width > slide_h * 1.2
                and shape.height < slide_h * 0.12
            )
            is_subtitle = (
                not is_title
                and shape.top < slide_h * 0.25
                and shape.width > slide_h * 0.8
                and shape.height > slide_h * 0.08
            )

            if is_label:
                slots["label"].append({"shape": shape, "area": area, "text": text, "top": shape.top})
            elif is_title:
                slots["title"].append({"shape": shape, "area": area, "top": shape.top})
            elif is_subtitle:
                slots["subtitle"].append({"shape": shape, "area": area, "top": shape.top})
            else:
                slots["content"].append({"shape": shape, "area": area, "top": shape.top})

        for key in slots:
            slots[key].sort(key=lambda s: s["top"])

        title_slot = slots["title"][0] if slots["title"] else None
        subtitle_slot = slots["subtitle"][0] if slots["subtitle"] else None
        content_slots = slots["content"]
        label_slots = slots["label"]

        if title_slot and title:
            _set_shape_text(title_slot["shape"], title)

        if subtitle_slot:
            _set_shape_text(subtitle_slot["shape"], subtitle)

        for idx, slot in enumerate(content_slots):
            if idx < len(modules):
                mod = modules[idx]
                text = mod["title"]
                for b in mod["bullets"][:3]:
                    text += "\n• " + b
                _set_shape_text(slot["shape"], text)
            else:
                _set_shape_text(slot["shape"], "")

        for idx, slot in enumerate(label_slots):
            orig = slot["text"]
            if orig.isdigit():
                pass
            else:
                _set_shape_text(slot["shape"], str(idx + 1))

    def _find_shape_by_id_or_name(self, slide, shape_id, shape_name):
        """Find a shape by ID or name."""
        # Try by shape_id first
        if shape_id is not None:
            for shape in slide.shapes:
                if getattr(shape, 'shape_id', None) == shape_id:
                    return shape
        # Fallback by name
        if shape_name:
            for shape in slide.shapes:
                if shape.name == shape_name:
                    return shape
        return None

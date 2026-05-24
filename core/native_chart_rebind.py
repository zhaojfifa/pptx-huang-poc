"""
Native chart rebind (PR-Q2E).

Updates native PowerPoint chart data IN PLACE using python-pptx `replace_data`, so charts
stay native/editable with their original style, axes, legend, colors and position intact.

- Only charts backed by an EMBEDDED workbook can be updated by `replace_data`; charts whose
  data is an external link raise `ValueError` (target-mode external) and are safely skipped.
- No matplotlib / SVG / PNG. No template/cloner redesign. Operates on an already-rendered
  output .pptx (the clone), so it cannot affect templates that contain no native charts
  (e.g. t5) — for those it is a no-op.
"""

import logging

from pptx import Presentation
from pptx.chart.data import CategoryChartData

logger = logging.getLogger(__name__)


def _make_chart_data(categories, series):
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values in series:
        cd.add_series(name, tuple(values))
    return cd


def rebind_native_charts(pptx_path: str, categories, series, only_slide_numbers=None) -> dict:
    """Rebind every embedded-workbook native chart in `pptx_path` to (categories, series).

    Args:
        pptx_path: path to a rendered .pptx (modified in place).
        categories: list of category labels (e.g. ["Q1","Q2","Q3","Q4"]).
        series: list of (series_name, values) tuples.
        only_slide_numbers: optional 1-based slide numbers to restrict rebind to (e.g. [7]).

    Returns a stats dict: charts_total / rebound / skipped_external / errors / slides.
    """
    stats = {"charts_total": 0, "rebound": 0, "skipped_external": 0, "errors": 0, "slides": []}
    prs = Presentation(pptx_path)
    changed = False
    for idx, slide in enumerate(prs.slides, start=1):
        if only_slide_numbers and idx not in only_slide_numbers:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            stats["charts_total"] += 1
            cd = _make_chart_data(categories, series)
            try:
                shape.chart.replace_data(cd)  # preserves style/axes/legend/colors/position
                stats["rebound"] += 1
                changed = True
                if idx not in stats["slides"]:
                    stats["slides"].append(idx)
            except ValueError as e:
                # external-link workbook (no embedded xlsx) → cannot rebind safely; skip.
                if "external" in str(e).lower():
                    stats["skipped_external"] += 1
                    logger.info(f"chart rebind skip (external link) on slide {idx}: {shape.name}")
                else:
                    stats["errors"] += 1
                    logger.warning(f"chart rebind error on slide {idx} {shape.name}: {e}")
            except Exception as e:
                stats["errors"] += 1
                logger.warning(f"chart rebind error on slide {idx} {shape.name}: {e}")
    if changed:
        prs.save(pptx_path)
    return stats

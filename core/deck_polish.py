"""
Deck polish (PR-Q2F) — conservative, post-render text cleanup on the generated .pptx.

Three small, safe passes (each logged, none deletes meaningful generated content):
  A. fix_agenda_slide   — fill agenda item slots with the REAL section titles, clear extras.
  B. strip_unsafe_numbering — remove leading section numbering ("2.1 ", "2. ") from titles,
       but never from 表/图/Table/Fig captions.
  C. clear_placeholder_residue — clear obvious unfilled template placeholders (环形饼图,
       类别&对应数量, 配图, 图例, 数据对应部门名称, xxxx, and bare "XX" fragments).

Everything is best-effort and reversible-by-regeneration; failures are swallowed per shape.
"""

import logging
import re

from pptx import Presentation

logger = logging.getLogger(__name__)

PLACEHOLDER_MARKERS = ["环形饼图", "类别&对应数量", "类别&数量", "配图", "图例",
                       "数据对应部门名称", "汇报内容概览占位", "xxxx", "XXXX"]
# leading section numbering on a TITLE: "2.1 ", "2.1.3 ", "2、", "2. " (NOT 表2-1 / 图2-1)
_MULTI_NUM = re.compile(r"^\s*\d+(?:\.\d+)+[\s、.．:：)）]*")
_SINGLE_NUM = re.compile(r"^\s*\d+[、.．)）]\s+")
_CAPTION = re.compile(r"^\s*(表|图|附表|附图|Table|Fig|Figure)", re.I)
_XX_ONLY = re.compile(r"^[Xx]{2,}$")
_XX_SUFFIX = re.compile(r"[:：]\s*[Xx]{2,}\s*$")


def _is_title_shape(shape) -> bool:
    name = (shape.name or "")
    return ("标题" in name) or ("title" in name.lower())


def _find_shape(slide, shape_id, name):
    """Find a shape by shape_id (preferred) then by exact name."""
    if shape_id is not None:
        for sh in slide.shapes:
            if getattr(sh, "shape_id", None) == shape_id:
                return sh
    if name:
        for sh in slide.shapes:
            if (sh.name or "") == name:
                return sh
    return None


def fix_cover_slide(prs, cover_title, cover_subtitle, log) -> None:
    """P1.6 (cover slide only): force the cover's visible title/subtitle to the confirmed-outline
    values. Targeted, never a renderer/global change:
      - title  → the visible title shape: prefer a shape whose name contains 标题/title; only if that
                 name heuristic fails, the conservative topmost text shape.
      - subtitle → only a shape whose name contains 副标题/subtitle (no fallback; skipped if absent).
    Writes ONLY those one/two shapes via _set_first_para; never blanks arbitrary shapes, never touches
    any other slide, table, or residue. Guarantees the rendered slide-1 title equals the outline title
    even when the analyzer mapped the blueprint title slot to a wrong/invisible shape."""
    if not cover_title and not cover_subtitle:
        return
    try:
        slide = list(prs.slides)[0]
    except IndexError:
        return
    text_shapes = [sh for sh in slide.shapes if sh.has_text_frame]
    if not text_shapes:
        return
    title_sh = next((sh for sh in text_shapes if _is_title_shape(sh)), None)
    if title_sh is None:
        title_sh = sorted(text_shapes, key=lambda s: (s.top or 0))[0]  # conservative fallback only
    if cover_title and title_sh is not None:
        _set_first_para(title_sh, cover_title)
        log["cover_set"].append({"shape": title_sh.name, "field": "title", "text": cover_title[:40]})
    if cover_subtitle:
        sub_sh = next((sh for sh in text_shapes if sh is not title_sh
                       and (("副标题" in (sh.name or "")) or ("subtitle" in (sh.name or "").lower()))), None)
        if sub_sh is not None:
            _set_first_para(sub_sh, cover_subtitle)
            log["cover_set"].append({"shape": sub_sh.name, "field": "subtitle", "text": cover_subtitle[:40]})


def fill_agenda_deterministic(prs, agenda_slide_number, section_titles, slot_refs, log) -> bool:
    """P1 (agenda-slide-only): overwrite the agenda page's content-slot shapes with the EXACT
    deduped body section_titles (in order), clearing extra agenda slots. Deterministic — does not
    depend on detecting a repeated 'generic' value — so the rendered 目录 matches 正文 section_title
    (the LLM's rephrasing/numbering of the agenda items is replaced). Returns True if applied.

    Scope guarantee: touches ONLY the agenda slide's content-slot shapes (identified by
    `slot_refs` = ordered (shape_id, name) from the agenda page blueprint). Never touches any other
    slide, never touches table cells, never cleans residue. If the template has fewer agenda slots
    than sections, only the leading sections are shown (template-side capacity; logged as
    agenda_overflow)."""
    if not agenda_slide_number or not section_titles or not slot_refs:
        return False
    try:
        slide = list(prs.slides)[agenda_slide_number - 1]
    except IndexError:
        return False
    applied = 0
    for i, (sid, sname) in enumerate(slot_refs):
        sh = _find_shape(slide, sid, sname)
        if sh is None or not sh.has_text_frame:
            continue
        if i < len(section_titles):
            _set_first_para(sh, section_titles[i])
            log["agenda_filled"].append({"slide": agenda_slide_number, "text": section_titles[i][:40]})
            applied += 1
        else:
            _set_first_para(sh, "")
            log["agenda_cleared"].append({"slide": agenda_slide_number})
    if len(section_titles) > len(slot_refs):
        log.setdefault("agenda_overflow", []).append(
            {"sections": len(section_titles), "slots": len(slot_refs)})
    return applied > 0


def fix_agenda_slide(prs, agenda_slide_number, agenda_items, log) -> None:
    """On the agenda slide, replace generic/duplicated item text with the real section
    titles in order; clear any leftover item slots. Conservative: only touches non-title
    text shapes whose text is the repeated generic value (the mis-filled item slots)."""
    if not agenda_slide_number or not agenda_items:
        return
    try:
        slide = list(prs.slides)[agenda_slide_number - 1]
    except IndexError:
        return
    # candidate item shapes: non-title, non-numeric, with text
    cand = []
    texts = {}
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t or _is_title_shape(sh):
            continue
        if t.isdigit() or re.fullmatch(r"[0-9０-９]{1,3}", t):
            continue  # keep numbering labels (01/02/03)
        cand.append(sh)
        texts[t] = texts.get(t, 0) + 1
    # the mis-filled item slots are the most-repeated generic value (count >= 2)
    generic = sorted([t for t, c in texts.items() if c >= 2], key=lambda t: -texts[t])
    if not generic:
        return  # nothing obviously wrong; do not touch
    generic_val = generic[0]
    item_shapes = [sh for sh in cand if sh.text_frame.text.strip() == generic_val]
    # order top-to-bottom, left-to-right
    item_shapes.sort(key=lambda s: (round((s.top or 0) / 360000), s.left or 0))
    for i, sh in enumerate(item_shapes):
        if i < len(agenda_items):
            _set_first_para(sh, agenda_items[i])
            log["agenda_filled"].append({"slide": agenda_slide_number, "text": agenda_items[i]})
        else:
            _set_first_para(sh, "")
            log["agenda_cleared"].append({"slide": agenda_slide_number})


def _set_first_para(shape, text: str) -> None:
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = text
    # drop extra paragraphs
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)


def strip_unsafe_numbering(prs, log) -> None:
    """Strip leading section numbering from short, title-like shapes (not 表/图 captions)."""
    for idx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            p = sh.text_frame.paragraphs[0] if sh.text_frame.paragraphs else None
            if not p or not p.runs:
                continue
            first = p.runs[0].text
            if not first or _CAPTION.match(first):
                continue
            title_like = _is_title_shape(sh) or len(sh.text_frame.text.strip()) <= 40
            if not title_like:
                continue
            new = _MULTI_NUM.sub("", first)
            if new == first:
                new = _SINGLE_NUM.sub("", first)
            if new != first and new.strip():
                p.runs[0].text = new
                log["numbering_stripped"].append({"slide": idx, "from": first[:30], "to": new[:30]})


def clear_placeholder_residue(prs, log) -> None:
    """Clear obvious unfilled template placeholders. Conservative: clears whole shape only
    when text contains a marker phrase or is a bare XX fragment; strips trailing ': XX'."""
    for idx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t:
                continue
            if any(m in t for m in PLACEHOLDER_MARKERS) or _XX_ONLY.match(t):
                _set_first_para(sh, "")
                log["placeholders_cleared"].append({"slide": idx, "text": t[:40]})
            elif _XX_SUFFIX.search(t):
                _set_first_para(sh, _XX_SUFFIX.sub("", t).rstrip("：: "))
                log["placeholders_cleared"].append({"slide": idx, "text": t[:40]})


def polish_deck(pptx_path: str, agenda_slide_number=None, agenda_items=None,
                agenda_slot_refs=None, cover_title=None, cover_subtitle=None) -> dict:
    log = {"agenda_filled": [], "agenda_cleared": [], "numbering_stripped": [],
           "placeholders_cleared": [], "cover_set": []}
    prs = Presentation(pptx_path)
    try:
        fix_cover_slide(prs, cover_title, cover_subtitle, log)   # P1.6 cover slide only
    except Exception as e:
        logger.warning(f"cover fix failed: {e}")
    try:
        # P1: when the agenda page's content-slot refs are known, deterministically fill the
        # agenda items with the exact deduped body section_titles (agenda slide only). Otherwise
        # fall back to the conservative repeated-generic fixer. No other slide is touched.
        done = False
        if agenda_slot_refs:
            done = fill_agenda_deterministic(prs, agenda_slide_number, agenda_items or [],
                                             agenda_slot_refs, log)
        if not done:
            fix_agenda_slide(prs, agenda_slide_number, agenda_items or [], log)
    except Exception as e:
        logger.warning(f"agenda fix failed: {e}")
    try:
        strip_unsafe_numbering(prs, log)
    except Exception as e:
        logger.warning(f"numbering strip failed: {e}")
    try:
        clear_placeholder_residue(prs, log)
    except Exception as e:
        logger.warning(f"placeholder cleanup failed: {e}")
    prs.save(pptx_path)
    log["summary"] = {k: len(v) for k, v in log.items() if isinstance(v, list)}
    return log

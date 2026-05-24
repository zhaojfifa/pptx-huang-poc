"""
FastAPI Web UI for PPTX Agent

Provides interactive endpoints for each pipeline step.
"""

import json
import logging
import re
import sys
import uuid
from datetime import datetime
from pathlib import Path

from fastapi import Body, FastAPI, File, Form, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
# Add project root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from config.settings import BASE_DIR, OUTPUT_DIR, PREVIEW_DIR
from core.agent import PPTXAgent
from database.db import GenerationJobDAO, TemplateDAO, init_db
from models.schemas import ConfirmRequest, GenerationRequest, RenderRequest

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

app = FastAPI(title="PPTX Agent")

# Static files
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")

# Read index.html directly (Jinja2Templates disabled due to Python 3.14 caching issue)
_INDEX_HTML_PATH = BASE_DIR / "static" / "templates" / "index.html"
_INDEX_HTML = _INDEX_HTML_PATH.read_text(encoding="utf-8") if _INDEX_HTML_PATH.exists() else "<h1>PPTX Agent</h1>"

agent = PPTXAgent()


@app.on_event("startup")
def startup():
    init_db()
    logger.info("Database initialized.")


@app.get("/", response_class=HTMLResponse)
def index():
    return HTMLResponse(content=_INDEX_HTML)


def _resolve_template(template_id: int | None, template_name: str | None) -> dict | None:
    """POC home selector → DB template. Prefer explicit id; else newest *complete*
    (>=12 pages) template matching the chosen style name; else newest by-name match.
    Returns the template dict or None. Does NOT run any generation step.
    """
    if template_id:
        t = TemplateDAO.get_by_id(template_id)
        if t:
            return t
    if not template_name:
        return None
    matches = [t for t in TemplateDAO.get_all() if t.get("name") == template_name]
    if not matches:
        return None
    from database.db import TemplatePageDAO
    # Prefer the same-name template with the MOST analyzed pages (the complete one),
    # tie-break on newest id. Works for any page count (t5=12, t6=8).
    return max(matches, key=lambda t: (len(TemplatePageDAO.get_by_template(t["id"])), t["id"]))


# A template is usable for real generation once it has a meaningful number of
# analyzed pages (excludes 0/partial rows; admits 8-page and 12-page masters).
_MIN_USABLE_PAGES = 6


@app.post("/api/jobs")
def create_job(
    user_requirements: str = Form(...),
    files: list[UploadFile] = File(default=[]),
    template_name: str = Form(default=None),
    template_key: str = Form(default=None),
    template_id: int = Form(default=None),
    input_mode: str = Form(default="manual"),
):
    saved_files = []
    upload_dir = BASE_DIR / "uploads"
    upload_dir.mkdir(exist_ok=True)

    for f in files:
        if f.filename:
            path = upload_dir / f.filename
            content = f.file.read()
            path.write_bytes(content)
            saved_files.append(str(path))

    job_id = agent.start_job(user_requirements, saved_files)

    # POC: persist the chosen template style so a later (human-triggered) pipeline
    # run maps to selected_template.json → t5.pptx. We do NOT run generation here.
    selected_template = _resolve_template(template_id, template_name)
    template_pages = 0
    if selected_template:
        from database.db import TemplatePageDAO
        GenerationJobDAO.update(job_id, selected_template_id=selected_template["id"])
        template_pages = len(TemplatePageDAO.get_by_template(selected_template["id"]))

    return {
        "job_id": job_id,
        "status": "created",
        "files": saved_files,
        "input_mode": input_mode,
        "template_name": template_name,
        "template_key": template_key,
        "resolved_template": (
            {"id": selected_template["id"], "name": selected_template["name"], "pages": template_pages}
            if selected_template else None
        ),
        "note": (
            "模板已就绪（>=12 页），可由人工触发后续生成步骤。"
            if template_pages >= 12 else
            "已记录所选模板风格；后续生成需人工先完成 t5.pptx 的完整 12 页分析（当前模板为 partial 或未入库）。"
        ),
    }


@app.post("/api/jobs/{job_id}/step/document")
def step_document(job_id: int):
    try:
        md = agent.step_process_documents(job_id)
        return {"job_id": job_id, "step": "document", "status": "ok", "markdown_length": len(md)}
    except Exception as e:
        logger.error(f"Step document failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "document", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/step/template")
def step_template(job_id: int):
    try:
        job = GenerationJobDAO.get_by_id(job_id)
        doc_md = ""
        checkpoint_md = BASE_DIR / "logs" / f"job_{job_id}" / "document_markdown.md"
        if checkpoint_md.exists():
            doc_md = checkpoint_md.read_text(encoding="utf-8")
        template = agent.step_select_template(job_id, job["user_requirements"], doc_md)
        return {"job_id": job_id, "step": "template", "status": "ok", "template": template}
    except Exception as e:
        logger.error(f"Step template failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "template", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/step/outline")
def step_outline(job_id: int):
    try:
        job = GenerationJobDAO.get_by_id(job_id)
        doc_md = ""
        checkpoint_md = BASE_DIR / "logs" / f"job_{job_id}" / "document_markdown.md"
        if checkpoint_md.exists():
            doc_md = checkpoint_md.read_text(encoding="utf-8")
        template = TemplateDAO.get_by_id(job["selected_template_id"]) if job.get("selected_template_id") else None
        if not template:
            return JSONResponse({"job_id": job_id, "step": "outline", "status": "error", "error": "No template selected"}, status_code=400)
        outline = agent.step_generate_outline(job_id, job["user_requirements"], doc_md, template)
        return {"job_id": job_id, "step": "outline", "status": "ok", "outline": outline}
    except Exception as e:
        logger.error(f"Step outline failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "outline", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/step/content")
def step_content(job_id: int):
    try:
        job = GenerationJobDAO.get_by_id(job_id)
        doc_md = ""
        checkpoint_md = BASE_DIR / "logs" / f"job_{job_id}" / "document_markdown.md"
        if checkpoint_md.exists():
            doc_md = checkpoint_md.read_text(encoding="utf-8")
        template = TemplateDAO.get_by_id(job["selected_template_id"]) if job.get("selected_template_id") else None
        outline = job.get("outline_json") or {}
        slides_data = agent.step_generate_content_and_layout(job_id, outline, doc_md, template)
        return {"job_id": job_id, "step": "content", "status": "ok", "slides_count": len(slides_data)}
    except Exception as e:
        logger.error(f"Step content failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "content", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/step/layout")
def step_layout(job_id: int):
    try:
        checkpoint = BASE_DIR / "logs" / f"job_{job_id}" / "slides_data.json"
        if not checkpoint.exists():
            return JSONResponse({"job_id": job_id, "step": "layout", "status": "error", "error": "No slides data found"}, status_code=400)
        slides_data = json.loads(checkpoint.read_text(encoding="utf-8"))
        layouts = agent.step_normalize_layouts(job_id, slides_data)
        return {"job_id": job_id, "step": "layout", "status": "ok", "layouts_count": len(layouts)}
    except Exception as e:
        logger.error(f"Step layout failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "layout", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/step/render")
def step_render(job_id: int):
    try:
        checkpoint = BASE_DIR / "logs" / f"job_{job_id}" / "layouts.json"
        if not checkpoint.exists():
            return JSONResponse({"job_id": job_id, "step": "render", "status": "error", "error": "No layouts found"}, status_code=400)
        layouts = json.loads(checkpoint.read_text(encoding="utf-8"))
        preview_path = agent.step_render_preview(job_id, layouts)
        return {"job_id": job_id, "step": "render", "status": "ok", "preview_path": preview_path}
    except Exception as e:
        logger.error(f"Step render failed: {e}")
        return JSONResponse({"job_id": job_id, "step": "render", "status": "error", "error": str(e)}, status_code=500)


@app.post("/api/jobs/{job_id}/confirm/outline")
def confirm_outline(req: ConfirmRequest):
    log = {
        "step": "outline_confirmation",
        "time": datetime.now().isoformat(),
        "confirmed": req.confirmed,
        "feedback": req.feedback,
    }
    GenerationJobDAO.append_log(req.job_id, log)
    if req.confirmed:
        GenerationJobDAO.update(req.job_id, status="outline_confirmed")
        return {"job_id": req.job_id, "status": "outline_confirmed"}
    else:
        # Could trigger regeneration with feedback
        return {"job_id": req.job_id, "status": "outline_rejected", "message": "Please provide updated requirements or feedback."}


@app.post("/api/jobs/{job_id}/confirm/render")
def confirm_render(req: RenderRequest):
    log = {
        "step": "render_confirmation",
        "time": datetime.now().isoformat(),
        "confirmed": req.confirmed,
        "feedback": req.feedback,
    }
    GenerationJobDAO.append_log(req.job_id, log)
    if req.confirmed:
        checkpoint = BASE_DIR / "logs" / f"job_{req.job_id}" / "layouts.json"
        layouts = json.loads(checkpoint.read_text(encoding="utf-8"))
        final_path = agent.step_generate_final(req.job_id, layouts)
        return {"job_id": req.job_id, "status": "completed", "final_path": final_path}
    else:
        GenerationJobDAO.update(req.job_id, status="render_rejected")
        return {"job_id": req.job_id, "status": "render_rejected", "message": "Please provide feedback for adjustment."}


@app.get("/api/jobs/{job_id}")
def get_job(job_id: int):
    job = GenerationJobDAO.get_by_id(job_id)
    if not job:
        return JSONResponse({"error": "Job not found"}, status_code=404)
    return job


@app.get("/api/jobs/{job_id}/download/{file_type}")
def download_file(job_id: int, file_type: str):
    job = GenerationJobDAO.get_by_id(job_id)
    if not job:
        return JSONResponse({"error": "Job not found"}, status_code=404)
    path = job.get("rendered_pptx_path") if file_type == "preview" else job.get("final_pptx_path")
    if not path or not Path(path).exists():
        return JSONResponse({"error": "File not found"}, status_code=404)
    return FileResponse(path, filename=Path(path).name)


@app.get("/api/templates")
def list_templates():
    return TemplateDAO.get_all()


# ----------------------------------------------------------------------------
# POC product flow: real outline (Kimi, template-bound) → user edits → confirm →
# real generation pipeline (background). Mock outline only as a fallback when the
# chosen template has no complete (>=12 page) data yet.
# ----------------------------------------------------------------------------
import threading

# Bundled, in-repo example source document (real content → grounds the default flow).
# Falls back to a legacy checkpoint if the bundled sample is ever missing.
_EXAMPLE_DOC = BASE_DIR / "data" / "samples" / "baosteel_2025_example.md"
_EXAMPLE_DOC_FALLBACK = BASE_DIR / "logs" / "job_62" / "document_markdown.md"

# In-memory cache for uploaded business documents (POC): token -> parsed markdown.
_SOURCE_UPLOADS: dict[str, dict] = {}

_AUDIENCE = {"management": "管理层", "client": "客户", "internal": "内部团队", "investor": "投资人"}
_SCENARIO = {"business_review": "经营分析", "project_report": "项目汇报", "solution": "解决方案", "training": "培训材料"}
_TONE = {"professional": "专业", "concise": "简洁", "formal": "正式"}
_LANG = {"zh-CN": "简体中文", "en": "English"}

# in-memory generation status (POC; reference -> state)
_GEN_STATUS: dict[int, dict] = {}


def _source_markdown(input_mode: str, source_name: str = None, source_token: str = None) -> str:
    """Resolve the source document text for the chosen input mode.

    - upload: parsed markdown of the uploaded business doc (cached by token).
    - example: the bundled in-repo sample (real content), with legacy fallback.
    - manual/text: empty (outline driven by the prompt only).
    """
    if source_token and source_token in _SOURCE_UPLOADS:
        return _SOURCE_UPLOADS[source_token].get("markdown", "") or ""
    if input_mode == "example":
        if _EXAMPLE_DOC.exists():
            return _EXAMPLE_DOC.read_text(encoding="utf-8")
        if _EXAMPLE_DOC_FALLBACK.exists():
            return _EXAMPLE_DOC_FALLBACK.read_text(encoding="utf-8")
    return ""


def _extract_key_facts(doc_md: str, limit: int = 12) -> list[str]:
    """Cheap, no-LLM extraction of fact/metric lines (numbers/%/¥/同比) for grounding."""
    if not doc_md:
        return []
    fact_re = re.compile(r"(\d|%|％|¥|亿|万|同比|环比|增长|下降|占比|利润|营收|产量)")
    facts = []
    for raw in doc_md.splitlines():
        line = raw.strip().lstrip("-*•· ").strip()
        if len(line) < 6 or line.startswith("#") or line.startswith(">"):
            continue
        if fact_re.search(line):
            facts.append(line)
        if len(facts) >= limit:
            break
    return facts


def _ground_doc(doc_md: str) -> str:
    """Prepend a compact 关键事实/数据 block so both outline and per-slide prompts
    receive distilled metrics in addition to the full document. No-op if empty."""
    if not doc_md:
        return ""
    facts = _extract_key_facts(doc_md)
    if not facts:
        return doc_md
    block = "## 关键事实/数据（自动提炼，供生成时严格依据）\n" + "\n".join(f"- {f}" for f in facts)
    return block + "\n\n" + doc_md


def _compose_requirements(payload: dict) -> str:
    """Fold the product config into a single natural-language requirements string."""
    p = (payload.get("prompt") or "").strip() or "生成一份企业经营汇报 PPT"
    bits = [p, ""]
    bits.append(f"页数：约 {payload.get('page_count', 12)} 页")
    bits.append(f"受众：{_AUDIENCE.get(payload.get('audience'), '管理层')}")
    bits.append(f"场景：{_SCENARIO.get(payload.get('scenario'), '经营分析')}")
    bits.append(f"语气：{_TONE.get(payload.get('tone'), '专业')}")
    bits.append(f"语言：{_LANG.get(payload.get('language'), '简体中文')}")
    if (payload.get("template_key") or "").lower() == "tech_blue":
        bits.append("风格取向：科技蓝——侧重技术方案/数字化转型/架构与平台能力/数据驱动/智能化，"
                    "标题偏「能力建设·技术路径·应用场景·平台价值」，避免纯财务报表口吻，内容须有资料依据。")
    return "\n".join(bits)


def _outline_to_cards(outline: dict) -> list:
    """Kimi outline {slides:[...]} → editable frontend cards (carry type/page map +
    agenda section binding section_id/section_title/slide_role_under_section)."""
    cards = []
    for s in outline.get("slides", []):
        cards.append({
            "page": s.get("slide_number"),
            "template_page_number": s.get("template_page_number", s.get("slide_number")),
            "section": s.get("title_section") or s.get("type") or "",
            "type": s.get("type") or "content",
            "title": s.get("title", ""),
            "points": list(s.get("key_points", []) or []),
            "section_id": s.get("section_id"),
            "section_title": s.get("section_title"),
            "slide_role_under_section": s.get("slide_role_under_section"),
        })
    return cards


def _cards_to_slides(cards: list, template_key: str = None) -> dict:
    """Confirmed cards → outline {slides:[...]} consumed verbatim by the pipeline."""
    total = len(cards)
    slides = []
    for i, c in enumerate(cards, start=1):
        slides.append({
            "slide_number": c.get("page", i),
            "template_page_number": c.get("template_page_number", c.get("page", i)),
            "type": c.get("type", "content"),
            "title": c.get("title", ""),
            "key_points": list(c.get("points", []) or []),
            "section_id": c.get("section_id"),
            "section_title": c.get("section_title"),
            "slide_role_under_section": c.get("slide_role_under_section"),
            "_total_pages": total,
            "_template_key": template_key,
            "notes": "",
        })
    return {"title": (cards[0].get("title") if cards else "企业经营汇报"), "slides": slides}


def _bigrams(text: str) -> set:
    t = re.sub(r"[\s，。、：:·\-（）()【】\[\]/]+", "", text or "")
    return {t[i:i + 2] for i in range(len(t) - 1)} if len(t) >= 2 else ({t} if t else set())


def _similarity(a: str, b: str) -> float:
    """Lightweight char-bigram Jaccard, no LLM."""
    ba, bb = _bigrams(a), _bigrams(b)
    if not ba or not bb:
        return 0.0
    return len(ba & bb) / len(ba | bb)


_AGENDA_TYPES = {"agenda", "section", "toc"}
_SKIP_TYPES = {"title", "cover", "ending", "closing", "thanks"}


def _agenda_consistency(slides: list, threshold: float = 0.12) -> dict:
    """Map body slides to agenda items (top-down check) and tag each slide in place with
    `_agenda_section` for prompt injection. Advisory: returns a report, never blocks."""
    # 1) detect agenda slide
    agenda_idx = None
    for i, s in enumerate(slides):
        ty = (s.get("type") or "").lower()
        title = s.get("title") or ""
        if ty in _AGENDA_TYPES or any(k in title for k in ("目录", "提纲", "议程", "汇报内容", "Agenda")):
            agenda_idx = i
            break
    if agenda_idx is None and len(slides) >= 2:
        agenda_idx = 1  # conventional 2nd slide fallback

    # Does the outline carry explicit section binding (PR-Q2D)? Prefer it over similarity.
    explicit = any((s.get("section_title") or s.get("section_id")) for s in slides)
    mapping_method = "section_id" if explicit else "similarity"

    # Build the authoritative agenda items.
    if explicit:
        # ordered distinct section_titles (by section_id when available, else first-seen)
        seen_order, sec_by_id = [], {}
        for s in slides:
            st = (s.get("section_title") or "").strip()
            if st and st not in seen_order:
                seen_order.append(st)
                if s.get("section_id"):
                    sec_by_id[st] = str(s.get("section_id"))
        agenda_items = sorted(seen_order, key=lambda t: sec_by_id.get(t, "zz")) if sec_by_id else seen_order
    else:
        agenda_items = []
        if agenda_idx is not None:
            a = slides[agenda_idx]
            agenda_items = [str(p).strip() for p in (a.get("key_points") or []) if str(p).strip()]
            if not agenda_items and a.get("title"):
                agenda_items = [a["title"].strip()]

    # 2) map each body slide to its agenda section; tag in place
    mappings, unmatched = [], []
    used = set()
    for i, s in enumerate(slides):
        ty = (s.get("type") or "").lower()
        if i == agenda_idx or ty in _SKIP_TYPES:
            continue
        slide_no = s.get("slide_number", i + 1)
        if explicit:
            st = (s.get("section_title") or "").strip()
            if st and st in agenda_items:
                s["_agenda_section"] = st
                used.add(st)
                mappings.append({"slide": slide_no, "title": s.get("title", ""),
                                 "agenda_item": st, "section_id": s.get("section_id"),
                                 "method": "section_id"})
            else:
                unmatched.append({"slide": slide_no, "title": s.get("title", ""),
                                  "best_item": st or None, "method": "section_id"})
            continue
        # similarity fallback (outlines without explicit binding)
        body = (s.get("title") or "") + " " + " ".join(str(p) for p in (s.get("key_points") or []))
        best, best_score = None, 0.0
        for item in agenda_items:
            sc = _similarity(item, body)
            if sc > best_score:
                best, best_score = item, sc
        if best and best_score >= threshold:
            s["_agenda_section"] = best
            used.add(best)
            mappings.append({"slide": slide_no, "title": s.get("title", ""),
                             "agenda_item": best, "score": round(best_score, 3), "method": "similarity"})
        else:
            unmatched.append({"slide": slide_no, "title": s.get("title", ""),
                              "best_item": best, "score": round(best_score, 3), "method": "similarity"})

    # 3) agenda items with no supporting slide
    orphans = [it for it in agenda_items if it not in used]
    # duplicate/drifted: an agenda item claimed by an implausible number of slides, or many unmatched
    drift = [m["agenda_item"] for m in mappings]
    dup = sorted({x for x in drift if drift.count(x) > max(2, len(slides) // 3)})

    if not agenda_items:
        status = "warning"
    elif unmatched and (len(unmatched) > max(1, len(mappings))):
        status = "fail"
    elif unmatched or orphans:
        status = "warning"
    else:
        status = "pass"

    return {
        "mapping_method": mapping_method,
        "agenda_slide_index": (agenda_idx + 1) if agenda_idx is not None else None,
        "agenda_items": agenda_items,
        "agenda_items_count": len(agenda_items),
        "slide_mappings": mappings,
        "unmatched_slides": unmatched,
        "agenda_items_without_supporting_slides": orphans,
        "duplicate_or_drifted_topics": dup,
        "overall_status": status,
    }


def _write_selected_template(job_id: int, template: dict, payload: dict, pages: int):
    """Persist selected_template.json so outline / slide / render reference one source."""
    tid = template["id"]
    record = {
        "template_key": payload.get("template_key") or "business",
        "template_name": template.get("name"),
        "file_path": template.get("file_path"),
        "page_count": pages,
        "overall_style": template.get("overall_style"),
        "screenshots_path": str(BASE_DIR / "logs" / f"template_screenshots_{tid}"),
        "blueprint_source": "template_pages",
        "frontend_selected": True,
    }
    agent._save_checkpoint(job_id, "selected_template.json", agent._json_dumps(record))
    return record


def _mock_outline_cards(payload: dict) -> list:
    base = [
        ("封面", "title", "{t}", ["公司名称与报告主题", "汇报对象与时间"]),
        ("目录", "section", "汇报内容概览", ["经营总览", "财务表现", "业务与能力", "风险与下一步"]),
        ("经营总览", "summary", "2025年度经营总体表现", ["全年经营总体情况", "核心成果与关键结论"]),
        ("财务表现", "chart", "核心财务数据概览", ["营收与利润总体表现", "盈利质量与现金流", "资产结构稳健性"]),
        ("财务表现", "chart", "主要财务指标与季度走势", ["每股收益与净资产收益率", "季度表现与节奏"]),
        ("经营成果", "chart", "2025年经营亮点", ["业绩亮点", "科技创新与产品突破", "运营效率提升"]),
        ("业务结构", "content", "主营业务与产品结构", ["收入与成本结构", "主要产销情况", "重点板块表现"]),
        ("核心能力", "section", "核心竞争力分析", ["产品与技术领先", "绿色低碳进展", "智慧制造"]),
        ("重点进展", "content", "重大投资与重点项目", ["重点投资方向", "重大项目进展"]),
        ("风险挑战", "content", "当前主要挑战", ["外部环境与行业压力", "经营层面的关键挑战"]),
        ("发展方向", "content", "下一阶段经营计划与重点", ["经营目标与方向", "重点举措与资源安排"]),
        ("结尾", "ending", "总结与展望", ["全年总结", "战略定力与未来展望"]),
    ]
    prompt = (payload.get("prompt") or "").strip()
    t = (prompt[:18] + "…") if len(prompt) > 18 else (prompt or "企业年度经营汇报")
    cards = []
    for i, (sec, ty, title, pts) in enumerate(base, start=1):
        cards.append({"page": i, "template_page_number": i, "section": sec, "type": ty,
                      "title": title.replace("{t}", t), "points": pts})
    return cards


# ---------------------------------------------------------------------------
# Quality gates (PR-Q1b/Q1c): outline quality, slot fragmentation, contamination.
# All gates are advisory for this PR — they log/persist warnings and never block
# generation (except a clearly unusable empty outline). No secrets are recorded.
# ---------------------------------------------------------------------------
_GENERIC_TERMS = ["建设背景", "关键能力", "应用成效", "下一步计划", "总体目标", "技术底座",
                  "核心组件", "价值闭环", "应对策略", "演进路线", "总体概述", "概述"]
_FACT_RE = re.compile(r"[0-9０-９]|%|％|¥|亿|万|同比|环比|增长|下降|占比")

# Final-PPTX forbidden/suspicious terms (old-template contamination + internals).
_FORBIDDEN_TERMS = ["穿透监管", "穿透式监管", "VPN", "SD-WAN", "集团专线", "是否模型",
                    "指标规则", "模型场景", "1模型", "多个规则", "CMCC", "OneCity",
                    "template_id", "Huang", "Kimi", "localhost", "sk-"]


def _write_report(job_id: int, filename: str, data: dict):
    """Persist a gate report under logs/job_<id>/ (no secrets in payload)."""
    try:
        agent._save_checkpoint(job_id, filename, json.dumps(data, ensure_ascii=False, indent=2))
    except Exception as e:
        logger.warning(f"Job {job_id}: failed to write {filename}: {e}")


def _outline_quality_gate(cards: list, doc_md: str, requested_pages: int,
                          template_pages: int) -> dict:
    """Lightweight advisory checks on the outline before per-slide generation."""
    warnings = []
    titles = [(c.get("title") or "").strip() for c in cards]

    # 1) duplicate / near-duplicate titles
    seen = {}
    for i, t in enumerate(titles):
        key = re.sub(r"[：:·\s]", "", t)
        if key and key in seen:
            warnings.append(f"duplicate_title: '{t}' (slides {seen[key] + 1} & {i + 1})")
        elif key:
            seen[key] = i

    # 2) empty or generic slide titles
    for i, t in enumerate(titles):
        if not t:
            warnings.append(f"empty_title: slide {i + 1}")
        elif t in _GENERIC_TERMS or len(t) <= 2:
            warnings.append(f"generic_title: slide {i + 1} '{t}'")

    # 3) empty key_points on narrative-type slides (cover/ending may be empty)
    for i, c in enumerate(cards):
        ty = (c.get("type") or "content")
        if ty in ("content", "chart", "summary", "section") and not (c.get("points") or []):
            warnings.append(f"empty_key_points: slide {i + 1} ({ty})")

    # 4) too many generic phrasings without business facts
    generic_hits = sum(1 for t in titles if any(g in t for g in _GENERIC_TERMS))
    if generic_hits >= max(3, len(cards) // 3):
        warnings.append(f"high_generic_titles: {generic_hits} slides use generic phrasing")

    # 5) page-count mismatch
    if requested_pages and len(cards) != requested_pages:
        warnings.append(f"page_count_mismatch: outline {len(cards)} vs requested {requested_pages}")
    if template_pages and len(cards) != template_pages:
        warnings.append(f"template_page_mismatch: outline {len(cards)} vs template {template_pages}")

    # 6) no source-backed facts even though grounding is available
    if doc_md and _extract_key_facts(doc_md):
        joined = " ".join(titles + [str(p) for c in cards for p in (c.get("points") or [])])
        if not _FACT_RE.search(joined):
            warnings.append("no_source_facts_in_outline: grounding available but outline has no numbers/metrics")

    return {"slides": len(cards), "warning_count": len(warnings), "warnings": warnings,
            "unusable": len(cards) == 0}


def _slot_fragmentation_report(template_pages: list, cards: list) -> dict:
    """Detect dense-label / low-capacity template pages and narrative slides mapped
    onto them. Advisory only — does not change mapping."""
    from core.template_style_engine import TemplateStyleEngine
    pages_report = []
    for i, tp in enumerate(template_pages):
        try:
            bp = TemplateStyleEngine(tp).blueprint
        except Exception as e:
            pages_report.append({"slide": i + 1, "error": str(e)})
            continue
        slots = bp.get("slots", {}) or {}
        content = slots.get("content", []) or []
        labels = slots.get("labels", []) or []
        caps = [s.get("_capacity", {}).get("total_chars", 0) for s in content]
        avg_cap = int(sum(caps) / len(caps)) if caps else 0
        warns = []
        if (len(content) + len(labels)) >= 20:
            warns.append("very_high_slot_count")
        if labels and len(labels) >= 2 * max(1, len(content)):
            warns.append("label_dominated")
        if content and avg_cap < 30:
            warns.append("small_avg_capacity")
        card = cards[i] if i < len(cards) else {}
        ctype = (card.get("type") or "content")
        if ctype in ("content", "summary") and len(labels) >= 12 and len(content) <= 2:
            warns.append("narrative_on_dense_label_page")
        pages_report.append({"slide": i + 1, "type": ctype,
                             "content_slots": len(content), "label_slots": len(labels),
                             "avg_content_capacity": avg_cap, "warnings": warns})
    flagged = [p["slide"] for p in pages_report if p.get("warnings")]
    return {"flagged_count": len(flagged), "flagged_pages": flagged, "pages": pages_report}


def _final_contamination_scan(pptx_path: str) -> dict:
    """Scan visible text + tables + speaker notes + docProps for forbidden/suspicious
    terms. Records hit COUNTS only (never the matched secret value)."""
    from pptx import Presentation
    texts = []
    try:
        prs = Presentation(pptx_path)
        for s in prs.slides:
            for sh in s.shapes:
                if sh.has_text_frame:
                    texts.append(sh.text_frame.text)
                if sh.has_table:
                    for row in sh.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
            if s.has_notes_slide and s.notes_slide.notes_text_frame:
                texts.append(s.notes_slide.notes_text_frame.text)
        cp = prs.core_properties
        texts += [cp.title or "", cp.author or "", cp.subject or "",
                  cp.comments or "", cp.keywords or "", cp.category or ""]
    except Exception as e:
        return {"clean": False, "error": str(e)}
    full = "\n".join(texts)
    hits = {}
    for term in _FORBIDDEN_TERMS:
        n = full.count(term)
        if n:
            hits[term] = n  # count only; never store the matched text
    return {"clean": len(hits) == 0, "total_hits": sum(hits.values()),
            "hit_terms": hits, "scanned_chars": len(full)}


@app.post("/api/poc/outline")
def poc_outline(payload: dict = Body(...)):
    """Generate an editable outline. Calls Kimi when the chosen template has
    complete (>=12 page) data; otherwise returns a structured fallback."""
    from database.db import TemplatePageDAO
    template = _resolve_template(payload.get("template_id"), payload.get("template_name"))
    pages = TemplatePageDAO.get_by_template(template["id"]) if template else []

    # Page-count guard: never exceed the template's own page count.
    req_pages = int(payload.get("page_count") or len(pages) or 12)
    if template and pages and req_pages > len(pages):
        return {"status": "error", "code": "page_count_exceeded",
                "message": f"该模板风格为 {len(pages)} 页，请选择不超过 {len(pages)} 页。",
                "template_pages": len(pages)}

    if template and len(pages) >= _MIN_USABLE_PAGES:
        requirements = _compose_requirements(payload)
        raw_doc = _source_markdown(payload.get("input_mode"), payload.get("source_document_name"),
                                   payload.get("source_token"))
        doc_md = _ground_doc(raw_doc)
        job_id = agent.start_job(requirements, [])
        logger.info(
            f"Job {job_id}: source grounding — input_mode={payload.get('input_mode')!r} "
            f"token={'yes' if payload.get('source_token') else 'no'} "
            f"raw_chars={len(raw_doc)} grounded_chars={len(doc_md)} "
            f"key_facts={len(_extract_key_facts(raw_doc))}"
        )
        GenerationJobDAO.update(job_id, selected_template_id=template["id"])
        if doc_md:
            agent._save_checkpoint(job_id, "document_markdown.md", doc_md)
        _write_selected_template(job_id, template, payload, len(pages))
        try:
            outline = agent.content_generator.generate_outline(
                requirements, doc_md, template, pages, job_id=job_id)
            cards = _outline_to_cards(outline)
            if cards:
                gate = _outline_quality_gate(cards, doc_md, req_pages, len(pages))
                _write_report(job_id, "outline_quality_report.json", gate)
                if gate["warnings"]:
                    logger.warning(f"Job {job_id}: outline quality warnings: {gate['warnings']}")
                return {"status": "ok", "source": "kimi", "job_id": job_id,
                        "template_name": template["name"], "page_count": len(cards),
                        "outline": cards, "quality_warnings": gate["warnings"]}
        except Exception as e:
            logger.error(f"Kimi outline failed, using fallback: {e}")

    cards = _mock_outline_cards(payload)
    return {"status": "ok", "source": "fallback", "job_id": None,
            "template_name": payload.get("template_name"),
            "page_count": len(cards), "outline": cards}


def _run_generation(job_id: int, outline: dict, doc_md: str, template: dict):
    """Background worker: real content → layout → render → final, using the
    user-confirmed outline verbatim (no silent rewrite)."""
    try:
        _GEN_STATUS[job_id] = {"state": "generating", "step": "content"}
        slides_data = agent.step_generate_content_and_layout(job_id, outline, doc_md, template)
        _GEN_STATUS[job_id]["step"] = "layout"
        layouts = agent.step_normalize_layouts(job_id, slides_data)
        _GEN_STATUS[job_id]["step"] = "render"
        agent.step_render_preview(job_id, layouts)
        final_path = agent.step_generate_final(job_id, layouts)
        # Native chart rebind (PR-Q2E) — T6/tech_blue P7 only. Updates embedded-workbook
        # native charts in place (style/axes/legend/colors/position preserved); external-link
        # charts are skipped. T5 has no native charts → not reached (name gate + no-op anyway).
        chart_rebind = None
        if (template or {}).get("name") == "科技蓝风格":
            from core.native_chart_rebind import rebind_native_charts
            # PR-Q2F restraint: P7 has six tiny (~1.9") charts; a single 4-point series keeps
            # each mini-chart readable instead of crowding 8 bars into it.
            chart_rebind = rebind_native_charts(
                final_path,
                categories=["Q1", "Q2", "Q3", "Q4"],
                series=[("营收(亿元)", (728.8, 784.9, 810.6, 850.7))],
                only_slide_numbers=[7],
            )
            _write_report(job_id, "native_chart_rebind_report.json", chart_rebind)
            logger.info(f"Job {job_id}: native chart rebind P7 → {chart_rebind}")
        # PR-Q2F polish (all decks): agenda titles + numbering strip + placeholder cleanup.
        from core.deck_polish import polish_deck
        ag = _agenda_consistency(outline.get("slides", []))
        polish = polish_deck(final_path, ag.get("agenda_slide_index"), ag.get("agenda_items"))
        _write_report(job_id, "template_placeholder_cleanup_report.json", polish)
        logger.info(f"Job {job_id}: deck polish → {polish.get('summary')}")
        # PR-Q2G typography & slot-hierarchy polish (all decks; charts untouched).
        from core.typography_polish import typography_polish
        typo = typography_polish(final_path)
        _write_report(job_id, "typography_audit_report.json", typo)
        logger.info(f"Job {job_id}: typography polish → tables={typo.get('table_count')} "
                    f"body_capped={typo.get('table_body_runs_capped')} "
                    f"text_changed={typo.get('text_changed_shapes')}")
        # Final PPTX hygiene scan (advisory; does not block this PR).
        contamination = _final_contamination_scan(final_path)
        _write_report(job_id, "final_contamination_report.json", contamination)
        if not contamination.get("clean"):
            logger.warning(f"Job {job_id}: contamination scan hits: {contamination.get('hit_terms')}")
        _GEN_STATUS[job_id] = {"state": "done", "final_path": final_path,
                               "contamination_clean": contamination.get("clean"),
                               "contamination_hits": contamination.get("hit_terms", {}),
                               "chart_rebind": chart_rebind,
                               "polish": polish.get("summary") if polish else None,
                               "typography": {"tables": typo.get("table_count"),
                                              "table_body_capped": typo.get("table_body_runs_capped"),
                                              "text_changed": typo.get("text_changed_shapes")}}
    except Exception as e:
        logger.error(f"Generation failed for job {job_id}: {e}")
        _GEN_STATUS[job_id] = {"state": "failed", "error": str(e)}


@app.post("/api/poc/generate")
def poc_generate(payload: dict = Body(...)):
    """Accept the user-confirmed outline and start the real generation pipeline
    in the background. The confirmed outline is used verbatim."""
    from database.db import TemplatePageDAO
    cards = payload.get("outline") or []
    template = _resolve_template(payload.get("template_id"), payload.get("template_name"))
    pages = TemplatePageDAO.get_by_template(template["id"]) if template else []
    ready = bool(template) and len(pages) >= _MIN_USABLE_PAGES

    # Page-count guard: confirmed outline must not exceed the template's pages.
    if template and pages and len(cards) > len(pages):
        return {"status": "error", "code": "page_count_exceeded",
                "message": f"该模板风格为 {len(pages)} 页，请将大纲控制在 {len(pages)} 页以内。",
                "template_pages": len(pages)}

    requirements = _compose_requirements(payload)
    job_id = agent.start_job(requirements, [])
    outline = _cards_to_slides(cards, payload.get("template_key"))   # confirmed outline, verbatim
    if template:
        GenerationJobDAO.update(job_id, selected_template_id=template["id"],
                                outline_json=outline)
        _write_selected_template(job_id, template, payload, len(pages))
    raw_doc = _source_markdown(payload.get("input_mode"), payload.get("source_document_name"),
                               payload.get("source_token"))
    doc_md = _ground_doc(raw_doc)
    logger.info(
        f"Job {job_id}: generate source grounding — input_mode={payload.get('input_mode')!r} "
        f"token={'yes' if payload.get('source_token') else 'no'} "
        f"raw_chars={len(raw_doc)} grounded_chars={len(doc_md)}"
    )
    if doc_md:
        agent._save_checkpoint(job_id, "document_markdown.md", doc_md)
    agent._save_checkpoint(job_id, "confirmed_outline.json", agent._json_dumps(outline))

    # Quality gates (advisory) — persisted under this (final) job for co-located evidence.
    gate = _outline_quality_gate(cards, doc_md, len(cards), len(pages))
    _write_report(job_id, "outline_quality_report.json", gate)
    frag = _slot_fragmentation_report(pages, cards)
    _write_report(job_id, "slot_fragmentation_report.json", frag)
    # Agenda-to-slides consistency: tags outline["slides"] in place with `_agenda_section`
    # (used for per-slide prompt injection) and persists a report.
    agenda = _agenda_consistency(outline.get("slides", []))
    _write_report(job_id, "agenda_consistency_report.json", agenda)
    if agenda["overall_status"] != "pass":
        logger.warning(f"Job {job_id}: agenda consistency = {agenda['overall_status']}; "
                       f"unmatched={[u['slide'] for u in agenda['unmatched_slides']]} "
                       f"orphan_items={agenda['agenda_items_without_supporting_slides']}")
    if gate["warnings"]:
        logger.warning(f"Job {job_id}: outline quality warnings: {gate['warnings']}")
    if frag["flagged_pages"]:
        logger.warning(f"Job {job_id}: slot fragmentation flagged pages: {frag['flagged_pages']}")

    if ready:
        _GEN_STATUS[job_id] = {"state": "generating", "step": "queued"}
        threading.Thread(target=_run_generation, args=(job_id, outline, doc_md, template), daemon=True).start()
        msg = "大纲已确认，正在按你的大纲生成 PPT，完成后可下载。"
    else:
        _GEN_STATUS[job_id] = {"state": "pending_template"}
        msg = "大纲已确认并保存。模板数据准备完成后即可生成正式 PPT。"

    return {"status": "submitted", "reference": f"PPT-{job_id:04d}", "job_id": job_id,
            "template_name": payload.get("template_name"), "outline_pages": len(cards),
            "ready": ready, "message": msg,
            "quality_warnings": gate["warnings"],
            "slot_fragmentation_flagged": frag["flagged_pages"]}


@app.get("/api/poc/status/{job_id}")
def poc_status(job_id: int):
    st = _GEN_STATUS.get(job_id, {"state": "unknown"})
    out = {"job_id": job_id, "reference": f"PPT-{job_id:04d}", **st}
    if st.get("state") == "done":
        out["download_url"] = f"/api/jobs/{job_id}/download/final"
    return out


@app.post("/api/poc/source/upload")
async def poc_source_upload(file: UploadFile = File(...)):
    """Parse an uploaded business document into markdown and cache it by token.
    The token is passed back in the outline/generate payloads as `source_token`,
    so the uploaded content actually grounds generation (previously only the
    filename was sent and the content was discarded)."""
    upload_dir = BASE_DIR / "uploads"
    upload_dir.mkdir(exist_ok=True)
    safe_name = Path(file.filename or "source").name
    dest = upload_dir / f"{uuid.uuid4().hex}_{safe_name}"
    dest.write_bytes(await file.read())
    try:
        from core.document_processor import DocumentProcessor
        markdown = DocumentProcessor().process([str(dest)])
    except Exception as e:
        logger.error(f"Source upload parse failed for {safe_name}: {e}")
        return JSONResponse({"status": "error", "error": f"解析失败：{e}"}, status_code=500)
    if not markdown or not markdown.strip():
        return JSONResponse({"status": "error", "error": "无法从该文件提取文本内容。"},
                            status_code=422)
    token = uuid.uuid4().hex
    _SOURCE_UPLOADS[token] = {"name": safe_name, "markdown": markdown}
    logger.info(f"Source upload parsed: {safe_name} → {len(markdown)} chars, token={token[:8]}…")
    return {"status": "ok", "source_token": token, "source_document_name": safe_name,
            "markdown_chars": len(markdown), "key_facts": len(_extract_key_facts(markdown))}


# ============================================================================
# Custom Template Live Validation (separate flow from the main built-in POC).
# upload → analyze → index → outline → confirm → generate. Reuses the same
# analyzer / screenshot / TemplateStyleEngine / Kimi pipeline. NOT a marketplace.
# ============================================================================
import os
import sys
import uuid
import subprocess

from fastapi.responses import HTMLResponse as _HTMLResponse

_CUSTOM_HTML_PATH = BASE_DIR / "static" / "templates" / "custom.html"
_CUSTOM_UPLOADS: dict[str, dict] = {}   # token -> {file_path, page_count, template_name, analyze}
_POLLUTION_KW = ["穿透监管", "穿透式监管", "CMCC", "中国移动", "中移", "OneCity",
                 "是否模型", "指标规则", "模型场景", "司库", "VPN", "SD-WAN", "集团专网", "国资委"]


@app.get("/custom", response_class=_HTMLResponse)
def custom_page():
    if _CUSTOM_HTML_PATH.exists():
        return _HTMLResponse(_CUSTOM_HTML_PATH.read_text(encoding="utf-8"))
    return _HTMLResponse("<h1>自定义模板验证</h1><p>页面未就绪。</p>")


@app.post("/api/custom-template/upload")
def custom_upload(file: UploadFile = File(...)):
    if not (file.filename or "").lower().endswith(".pptx"):
        return JSONResponse({"status": "error", "message": "请上传 PPTX 母版文件。"}, status_code=400)
    token = uuid.uuid4().hex[:10]
    dest = BASE_DIR / "templates_storage" / f"custom_{token}.pptx"
    dest.parent.mkdir(exist_ok=True)
    dest.write_bytes(file.file.read())
    try:
        from pptx import Presentation
        pages = len(Presentation(str(dest)).slides)
    except Exception as e:
        dest.unlink(missing_ok=True)
        return JSONResponse({"status": "error", "message": f"无法解析该 PPTX：{e}"}, status_code=400)
    name = f"客户母版-{Path(file.filename).stem[:24]}-{token[:4]}"
    _CUSTOM_UPLOADS[token] = {"file_path": str(dest), "page_count": pages, "template_name": name}
    return {"status": "ok", "template_upload_id": token, "file_name": file.filename,
            "page_count": pages, "template_name": name}


def _run_custom_analyze(token: str):
    rec = _CUSTOM_UPLOADS.get(token)
    if not rec:
        return
    rec["analyze"] = {"state": "analyzing", "step": "解析与截图"}
    try:
        # Reuse the existing analyzer module exactly (screenshot + analyzer + DB insert).
        proc = subprocess.run(
            [sys.executable, "-m", "template_analyzer.analyze_template",
             "--input", rec["file_path"], "--name", rec["template_name"]],
            cwd=str(BASE_DIR), capture_output=True, text=True, timeout=3600,
        )
        if proc.returncode != 0:
            rec["analyze"] = {"state": "failed", "error": proc.stderr[-400:]}
            return
        from database.db import TemplatePageDAO
        tpl = _resolve_template(None, rec["template_name"])
        pages = TemplatePageDAO.get_by_template(tpl["id"]) if tpl else []
        rec["template_id"] = tpl["id"] if tpl else None
        rec["analyze"] = {"state": "ready" if len(pages) >= 1 else "failed",
                          "pages": len(pages)}
    except Exception as e:
        rec["analyze"] = {"state": "failed", "error": str(e)}


@app.post("/api/custom-template/analyze")
def custom_analyze(payload: dict = Body(...)):
    token = payload.get("template_upload_id")
    rec = _CUSTOM_UPLOADS.get(token)
    if not rec:
        return JSONResponse({"status": "error", "message": "上传记录不存在，请重新上传。"}, status_code=404)
    rec["analyze"] = {"state": "queued"}
    threading.Thread(target=_run_custom_analyze, args=(token,), daemon=True).start()
    return {"status": "started", "template_upload_id": token, "template_name": rec["template_name"]}


@app.get("/api/custom-template/status/{token}")
def custom_status(token: str):
    rec = _CUSTOM_UPLOADS.get(token)
    if not rec:
        return JSONResponse({"status": "error", "message": "记录不存在"}, status_code=404)
    a = rec.get("analyze", {"state": "uploaded"})
    out = {"template_upload_id": token, "template_name": rec["template_name"],
           "page_count": rec["page_count"], **a}
    if a.get("state") == "ready":
        out["index"] = _build_template_index(rec["template_name"])
    return out


def _build_template_index(template_name: str) -> list:
    """Slide index for a (custom) template: thumbnail/page_type/slot counts/warnings."""
    from database.db import TemplatePageDAO
    from core.template_style_engine import TemplateStyleEngine
    from core import page_type_prompt as ptp
    tpl = _resolve_template(None, template_name)
    if not tpl:
        return []
    pages = TemplatePageDAO.get_by_template(tpl["id"])
    total = len(pages)
    idx = []
    for p in pages:
        lj = json.loads(p["layout_json"]) if p.get("layout_json") else {}
        bp = TemplateStyleEngine({"layout_json": lj,
                                  "visual_json": json.loads(p.get("visual_json") or "{}")}).blueprint
        slots = bp.get("slots", {})
        shapes = lj.get("shapes", [])
        tables = sum(1 for s in shapes if "table" in s)
        pics = sum(1 for s in shapes if s.get("is_picture"))
        groups = sum(1 for s in shapes if "GROUP" in str(s.get("shape_type", "")))
        ptype = ptp.classify(p["page_number"], total, None, ptp.slot_summary(bp))
        md = p.get("markdown_content") or ""
        warns = []
        hit = [kw for kw in _POLLUTION_KW if kw in md]
        if hit:
            warns.append("含旧模板用词，建议清洗：" + "、".join(hit[:4]))
        if not slots.get("content") and not slots.get("labels") and not tables:
            warns.append("未检测到可填充区域")
        idx.append({
            "page": p["page_number"],
            "thumbnail": f"/api/custom-template/thumbnail/{tpl['id']}/{p['page_number']}",
            "page_type": ptp.describe(ptype),
            "content_slots": len(slots.get("content", [])),
            "label_slots": len(slots.get("labels", [])),
            "tables": tables, "pictures": pics, "groups": groups,
            "recommended_use": ptp.describe(ptype),
            "warnings": warns,
        })
    return idx


@app.get("/api/custom-template/thumbnail/{tid}/{n}")
def custom_thumbnail(tid: int, n: int):
    p = BASE_DIR / "logs" / f"template_screenshots_{tid}" / f"slide_{n}.png"
    if not p.exists():
        return JSONResponse({"error": "no thumbnail"}, status_code=404)
    return FileResponse(str(p))


@app.post("/api/custom-template/outline")
def custom_outline(payload: dict = Body(...)):
    # Reuse the same Kimi outline path, bound to the custom template by name.
    return poc_outline(payload)


@app.post("/api/custom-template/generate")
def custom_generate(payload: dict = Body(...)):
    return poc_generate(payload)

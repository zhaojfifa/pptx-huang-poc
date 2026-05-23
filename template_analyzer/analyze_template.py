#!/usr/bin/env python3
"""
Template Analyzer Script

Usage:
    python -m template_analyzer.analyze_template --input template.pptx --name "My Template"
"""

import argparse
import json
import logging
import sys
from pathlib import Path

# Add parent to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt

from config.settings import LOGS_DIR, TEMPLATES_DIR
from core.template_style_engine import TemplateStyleEngine
from database.db import TemplateDAO, TemplatePageDAO, init_db
from skills.skill_llm.llm_skill import LLMSkill
from skills.skill_markitdown.markitdown_skill import MarkItDownSkill

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(LOGS_DIR / "template_analyzer.log", encoding="utf-8"),
    ],
)
logger = logging.getLogger(__name__)


def extract_shape_info(shape):
    """Extract detailed shape properties."""
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "rotation": shape.rotation,
    }

    # Text frame
    if shape.has_text_frame:
        tf = shape.text_frame
        paragraphs = []
        for para in tf.paragraphs:
            para_info = {
                "text": para.text,
                "alignment": str(para.alignment),
                "level": para.level,
                "runs": [],
            }
            # Extract paragraph-level font as fallback (important for placeholders with no runs)
            para_font = {}
            if para.font:
                para_font = {
                    "font_name": para.font.name,
                    "font_size": para.font.size.pt if para.font.size else None,
                    "bold": para.font.bold,
                    "italic": para.font.italic,
                    "underline": para.font.underline,
                }
                try:
                    para_font["color"] = str(para.font.color.rgb)
                except Exception:
                    para_font["color"] = None
            # Extract line spacing from paragraph XML
            line_spacing = None
            try:
                from pptx.oxml.ns import qn
                pPr = para._p.find(qn('a:pPr'))
                if pPr is not None:
                    lnSpc = pPr.find(qn('a:lnSpc'))
                    if lnSpc is not None:
                        spcPct = lnSpc.find(qn('a:spcPct'))
                        if spcPct is not None:
                            val = spcPct.get('val')
                            if val:
                                line_spacing = int(val) / 100000
                        spcPts = lnSpc.find(qn('a:spcPts'))
                        if spcPts is not None:
                            val = spcPts.get('val')
                            if val:
                                line_spacing = int(val) / 100  # convert to pt-like value
            except Exception:
                pass
            para_info["line_spacing"] = line_spacing

            para_info["font"] = para_font

            for run in para.runs:
                run_info = {"text": run.text}
                if run.font:
                    run_info["font_name"] = run.font.name
                    run_info["font_size"] = run.font.size.pt if run.font.size else None
                    run_info["bold"] = run.font.bold
                    run_info["italic"] = run.font.italic
                    run_info["underline"] = run.font.underline
                    try:
                        run_info["color"] = str(run.font.color.rgb)
                    except Exception:
                        run_info["color"] = None
                para_info["runs"].append(run_info)
            paragraphs.append(para_info)
        info["paragraphs"] = paragraphs
        info["word_wrap"] = tf.word_wrap
        info["auto_size"] = str(tf.auto_size)

    # Fill — guard whole block: Picture/Connector/GraphicFrame(table) have no .fill
    # and accessing the property itself raises AttributeError (which would otherwise
    # drop the entire shape, including tables and pictures).
    try:
        if shape.fill:
            info["fill_type"] = str(shape.fill.type)
            if shape.fill.type == 1:  # SOLID
                try:
                    info["fill_color"] = str(shape.fill.fore_color.rgb)
                except Exception:
                    info["fill_color"] = None
    except Exception:
        pass

    # Line — same guard (Picture/GraphicFrame have no .line)
    try:
        if shape.line:
            try:
                info["line_color"] = str(shape.line.color.rgb)
            except Exception:
                info["line_color"] = None
            info["line_width"] = shape.line.width.pt if shape.line.width else None
    except Exception:
        pass

    # Picture
    if shape.shape_type == 13:  # PICTURE
        info["is_picture"] = True

    # Table
    if shape.has_table:
        table = shape.table
        from pptx.oxml.ns import qn
        rows_data = []
        headers = []
        cell_styles = {}
        for r_idx, row in enumerate(table.rows):
            row_cells = []
            for c_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                row_cells.append(text)
                # Extract cell style once per unique style pattern
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                fill_color = None
                if tcPr is not None:
                    solidFill = tcPr.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgb = solidFill.find(qn('a:srgbClr'))
                        if srgb is not None:
                            fill_color = srgb.get('val')
                        else:
                            scheme = solidFill.find(qn('a:schemeClr'))
                            if scheme is not None:
                                fill_color = f"scheme:{scheme.get('val')}"
                # Font info from first paragraph/run
                font_size = None
                bold = None
                font_color = None
                if cell.text_frame.paragraphs:
                    p = cell.text_frame.paragraphs[0]
                    if p.runs:
                        r = p.runs[0]
                        font_size = r.font.size.pt if r.font.size else None
                        bold = r.font.bold
                        try:
                            font_color = str(r.font.color.rgb)
                        except Exception:
                            pass
                style_key = f"{fill_color}|{font_color}|{font_size}|{bold}"
                if style_key not in cell_styles:
                    cell_styles[style_key] = {
                        "fill_color": fill_color,
                        "font_color": font_color,
                        "font_size": font_size,
                        "bold": bold,
                    }
            if r_idx == 0:
                headers = row_cells
            else:
                rows_data.append(row_cells)
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "headers": headers,
            "sample_rows": rows_data[:3],
            "cell_styles": list(cell_styles.values())[:5],
        }

    return info


def extract_layout_json(prs: Presentation, slide_idx: int) -> dict:
    slide = prs.slides[slide_idx]
    shapes_info = []
    for shape in slide.shapes:
        try:
            shapes_info.append(extract_shape_info(shape))
        except Exception as e:
            logger.warning(f"Failed to extract shape {shape.name}: {e}")

    layout = {
        "slide_index": slide_idx,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "shapes": shapes_info,
        "slide_layout_name": slide.slide_layout.name if slide.slide_layout else None,
    }
    return layout


def split_markdown_by_slide(full_md: str, num_slides: int) -> list:
    """Split markitdown output into per-slide markdown.

    markitdown emits `<!-- Slide number: N -->` markers per slide. We segment on
    those so each page's markdown belongs to that page (fixes the prior crude
    even-line-split that drifted badly when shape density varied across pages).
    Falls back to even-line distribution only if no markers are present.
    """
    import re
    pages = ["" for _ in range(num_slides)]
    if not full_md:
        return pages

    marker = re.compile(r"<!--\s*Slide number:\s*(\d+)\s*-->", re.IGNORECASE)
    matches = list(marker.finditer(full_md))
    if matches:
        for i, m in enumerate(matches):
            n = int(m.group(1))
            start = m.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(full_md)
            if 1 <= n <= num_slides:
                pages[n - 1] = full_md[start:end].strip()
        return pages

    # Fallback: even-line distribution (legacy behavior)
    lines = full_md.split("\n")
    chunk = max(1, len(lines) // num_slides)
    for idx in range(num_slides):
        s = idx * chunk
        e = s + chunk if idx < num_slides - 1 else len(lines)
        pages[idx] = "\n".join(lines[s:e])
    return pages


def analyze_with_llm(llm: LLMSkill, markdown_text: str, layout_json: dict, page_number: int) -> dict:
    system = "You are a PPT design analyst. Analyze the visual design of a PPT page and output JSON only."
    prompt = f"""
Analyze the following PPT page (page {page_number}).

Markdown content:
```
{markdown_text[:2000]}
```

Layout data (shapes count: {len(layout_json.get('shapes', []))}):
```json
{json.dumps(layout_json, ensure_ascii=False, indent=2)[:3000]}
```

Please analyze and return a JSON object with these fields:
- "visual_summary": Brief description of the visual style and composition
- "color_scheme": List of dominant colors used
- "typography_style": Description of font usage and hierarchy
- "layout_pattern": Description of how elements are arranged
- "design_elements": List of notable design elements (shapes, images, icons, etc.)
- "mood": The overall mood/feeling of the design (professional, casual, modern, classic, etc.)

Return valid JSON only.
"""
    return llm.chat_structured(prompt, system=system)


def analyze_with_vision(llm: LLMSkill, image_path: str, page_number: int) -> dict:
    """Analyze a PPT template page screenshot using multimodal LLM vision."""
    system = "You are a PPT design analyst. Analyze the visual design of a PPT page from its screenshot. Output JSON only."
    prompt = f"""
Analyze this PPT template page screenshot (page {page_number}).

Focus on visual elements that are hard to extract programmatically:
- Exact color values and gradients (including gradient directions and stops)
- Shadow effects (drop shadows, inner shadows, glow effects)
- Transparency / opacity levels of shapes
- Image textures, patterns, or background images
- Border styles (rounded corners, bevel, 3D effects)
- Icon styles (flat, outline, filled, 3D)
- Spacing, padding, and alignment precision
- Animation cues (if visually apparent)

Return a JSON object with these fields:
{{
    "visual_summary": "Brief visual composition description",
    "color_scheme": {{
        "primary": "hex or description",
        "secondary": "hex or description",
        "accent": "hex or description",
        "background": "hex or description",
        "gradient_description": "Describe any gradients: direction, colors, stops"
    }},
    "shadow_effects": "Describe shadows: color, blur, offset, direction",
    "border_styles": "Describe borders: width, color, rounded corners, bevel",
    "icon_style": "Describe icon appearance",
    "texture_or_pattern": "Any textures, patterns, or background images",
    "transparency": "Any transparent or semi-transparent elements",
    "spacing_alignment": "Describe spacing, padding, alignment",
    "overall_mood": "mood/feeling of the design",
    "visual_notes": "Any other notable visual details"
}}

Return valid JSON only.
"""
    raw = llm.chat_with_image(prompt, image_path, system=system, json_mode=True)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        logger.error(f"Failed to parse vision LLM JSON response: {raw}")
        return {"raw_response": raw, "parse_error": True}


def _build_slot_analysis(blueprint: dict) -> str:
    """Build a human-readable slot analysis from blueprint for LLM prompt."""
    rows = blueprint["slots"].get("_rows", [])
    lines = []
    if not rows:
        contents = blueprint["slots"].get("content", [])
        if contents:
            rows = [contents]
        else:
            return "No content slots detected."

    for row_idx, row in enumerate(rows):
        lines.append(f"Row {row_idx + 1} ({len(row)} slots):")
        for slot in row:
            cap = slot.get("_capacity", {})
            total = cap.get("total_chars", 40)
            lines_count = cap.get("max_lines", 1)
            cpl = cap.get("chars_per_line", 20)
            w = slot.get("width", 0)
            h = slot.get("height", 0)
            if total > 80 and lines_count >= 2:
                role = "main_content"
                desc = f"main content area, multi-line, capacity ~{total} chars / {lines_count} lines"
            elif w > h * 2:
                role = "horizontal_label"
                desc = f"horizontal label/caption, single-line, capacity ~{total} chars"
            else:
                role = "label_or_extract"
                desc = f"label or concise extract, capacity ~{total} chars / {lines_count} line(s)"
            lines.append(f"  - {slot.get('name', 'unknown')}: {desc}")

    # Add subtitle info if present
    subtitle = blueprint["slots"].get("subtitle")
    if subtitle:
        cap = subtitle.get("_capacity", {})
        total = cap.get("total_chars", 120)
        lines_count = cap.get("max_lines", 2)
        lines.append(f"Subtitle slot: capacity ~{total} chars / {lines_count} lines")

    return "\n".join(lines)


def analyze_generation_hints(llm: LLMSkill, layout_json: dict, visual_json: dict, blueprint: dict, page_number: int) -> dict:
    """Analyze template page layout and generate content generation constraints."""
    system = "You are a PPT content generation strategist. Generate precise Chinese prompt constraints for an LLM content writer. Output JSON only."

    slot_analysis = _build_slot_analysis(blueprint)
    visual_summary = (visual_json or {}).get("visual_summary", "")
    layout_pattern = (visual_json or {}).get("layout_pattern", "")
    slide_w = layout_json.get("slide_width", 9144000)
    slide_h = layout_json.get("slide_height", 6858000)
    total_slots = blueprint.get("content_slots_count", 0)

    prompt = f"""Analyze the following PPT template page and generate content generation constraint hints.

Template page: {page_number}
Slide dimensions: {slide_w} x {slide_h} EMU
Total content slots: {total_slots}
Visual summary: {visual_summary}
Layout pattern: {layout_pattern}

Slot analysis (grouped by visual rows):
{slot_analysis}

Based on the actual layout above, generate the following fields IN CHINESE. Be specific and adapt to the actual layout — do NOT use generic descriptions if the layout has unique characteristics (e.g. card-based, column-based, top-bottom split, grid, etc.):

1. "layout_description": Describe how content slots are spatially organized, what visual relationship adjacent slots have, and define each slot role with its expected content format.
   - If rows exist: explain that same-row slots are visually adjacent and should have logically related content.
   - Define roles based on actual capacity: large multi-line slots need titles + bullets; small single-line slots need concise labels only.
   - Adapt to actual layout pattern (cards, columns, rows, freeform, etc.).

2. "content_relevance": Rules about how content in adjacent / same-row / related slots should thematically relate to each other. Give concrete examples when possible.

3. "capacity_constraints": Rules about title length, bullet count, and total characters PER SLOT ROLE. Base this strictly on the capacity numbers provided above.

4. "module_title_rule": Rule about module title length (e.g. "每个模块标题: 15个中文字以内" or different if layout suggests otherwise).

5. "format_example": A short markdown example showing the expected ### module format for this specific layout, covering the different slot roles present.

Return valid JSON only:
{{
    "layout_description": "...",
    "content_relevance": "...",
    "capacity_constraints": "...",
    "module_title_rule": "...",
    "format_example": "..."
}}
"""
    try:
        result = llm.chat_structured(prompt, system=system, enable_thinking=False)
        # Validate expected keys
        expected_keys = {"layout_description", "content_relevance", "capacity_constraints", "module_title_rule", "format_example"}
        if not expected_keys.issubset(result.keys()):
            missing = expected_keys - result.keys()
            logger.warning(f"Generation hints missing keys {missing} for page {page_number}. Filling defaults.")
            for key in missing:
                result[key] = ""
        return result
    except Exception as e:
        logger.error(f"Failed to generate generation_hints for page {page_number}: {e}")
        return {
            "layout_description": "",
            "content_relevance": "",
            "capacity_constraints": "",
            "module_title_rule": "",
            "format_example": "",
            "error": str(e),
        }


def analyze_overall_style(llm: LLMSkill, all_pages_markdown: list, all_pages_layout: list) -> dict:
    system = "You are a PPT design analyst. Analyze the overall style of a PPT template and output JSON only."
    combined_md = "\n---\n".join([f"Page {i+1}:\n{md[:800]}" for i, md in enumerate(all_pages_markdown)])
    prompt = f"""
Analyze the overall style of this PPT template based on all its pages.

Pages markdown:
```
{combined_md}
```

Each page has approximately {len(all_pages_layout[0].get('shapes', [])) if all_pages_layout else 0} shapes.
Slide dimensions: width={all_pages_layout[0].get('slide_width') if all_pages_layout else 'unknown'}, height={all_pages_layout[0].get('slide_height') if all_pages_layout else 'unknown'}

Please return a JSON object with:
- "style_name": A name for this template style
- "description": Overall description of the template style
- "color_palette": Primary and secondary colors
- "font_recommendations": Recommended fonts for title, subtitle, body
- "typical_layouts": Description of typical page layouts used
- "target_audience": Who this template is best for
- "content_suitability": What types of content work best
- "design_keywords": List of keywords describing the style

Return valid JSON only.
"""
    return llm.chat_structured(prompt, system=system)


def main():
    parser = argparse.ArgumentParser(description="Analyze PPT template and store in database")
    parser.add_argument("--input", required=True, help="Path to .pptx template file")
    parser.add_argument("--name", required=True, help="Template name")
    args = parser.parse_args()

    init_db()

    input_path = Path(args.input)
    if not input_path.exists():
        logger.error(f"Input file not found: {input_path}")
        sys.exit(1)

    # Copy template to storage
    TEMPLATES_DIR.mkdir(exist_ok=True)
    stored_path = TEMPLATES_DIR / input_path.name
    import shutil
    shutil.copy(str(input_path), str(stored_path))

    logger.info(f"Analyzing template: {input_path}")

    # Extract markdown
    md_skill = MarkItDownSkill()
    try:
        full_md = md_skill.convert(str(input_path))
    except Exception as e:
        logger.error(f"markitdown failed: {e}")
        full_md = ""

    # Extract per-page layout
    prs = Presentation(str(input_path))
    all_pages_md = []
    all_pages_layout = []

    llm = LLMSkill()

    # Per-slide markdown via slide-number markers (fixes crude even-split drift)
    page_md_list = split_markdown_by_slide(full_md, len(prs.slides))

    for idx in range(len(prs.slides)):
        logger.info(f"Processing slide {idx + 1}/{len(prs.slides)}...")
        layout = extract_layout_json(prs, idx)
        all_pages_layout.append(layout)
        all_pages_md.append(page_md_list[idx])

    # Analyze overall style
    logger.info("Analyzing overall style with LLM...")
    overall_style = analyze_overall_style(llm, all_pages_md, all_pages_layout)

    # Save template first to get template_id
    template_id = TemplateDAO.create(args.name, str(stored_path), overall_style)
    logger.info(f"Template created with ID: {template_id}")

    # Export screenshots for vision analysis (after template_id is known)
    screenshot_dir = LOGS_DIR / f"template_screenshots_{template_id}"
    screenshot_paths = []
    try:
        from skills.skill_ppt_screenshot.ppt_screenshot_skill import PPTScreenshotSkill
        screenshot_skill = PPTScreenshotSkill(width=1920, height=1080)
        screenshot_paths = screenshot_skill.export_slides(str(stored_path), str(screenshot_dir))
        logger.info(f"Exported {len(screenshot_paths)} screenshots to {screenshot_dir}")
    except Exception as e:
        logger.warning(f"Screenshot export failed: {e}. Vision analysis will be skipped.")

    # Analyze and save each page
    for idx, (page_md, layout) in enumerate(zip(all_pages_md, all_pages_layout)):
        logger.info(f"Analyzing page {idx + 1}...")

        # 1. Text-based analysis
        try:
            visual_json = analyze_with_llm(llm, page_md, layout, idx + 1)
        except Exception as e:
            logger.error(f"LLM text analysis failed for page {idx+1}: {e}")
            visual_json = {"error": str(e)}

        # 2. Vision-based analysis (if screenshot available)
        if idx < len(screenshot_paths):
            try:
                vision_json = analyze_with_vision(llm, screenshot_paths[idx], idx + 1)
                # Merge vision insights into visual_json
                visual_json["vision_analysis"] = vision_json
                visual_json["has_screenshot"] = True
                logger.info(f"Vision analysis completed for page {idx + 1}")
            except Exception as e:
                logger.warning(f"Vision analysis failed for page {idx+1}: {e}")
                visual_json["vision_analysis"] = {"error": str(e)}
                visual_json["has_screenshot"] = False
        else:
            visual_json["has_screenshot"] = False

        # 3. Build blueprint and generate per-page content generation hints
        try:
            template_page_data = {"layout_json": layout, "visual_json": visual_json}
            engine = TemplateStyleEngine(template_page_data)
            blueprint = engine.blueprint
            generation_hints = analyze_generation_hints(llm, layout, visual_json, blueprint, idx + 1)
            logger.info(f"Generation hints created for page {idx + 1}")
        except Exception as e:
            logger.error(f"Failed to build blueprint/generation_hints for page {idx+1}: {e}")
            generation_hints = None

        TemplatePageDAO.create(
            template_id=template_id,
            page_number=idx + 1,
            markdown_content=page_md,
            layout_json=layout,
            visual_json=visual_json,
            generation_hints=generation_hints,
        )
        logger.info(f"Page {idx + 1} saved.")

    logger.info("Template analysis complete.")
    print(f"Template ID: {template_id}")


if __name__ == "__main__":
    main()

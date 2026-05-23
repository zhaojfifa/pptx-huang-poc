import logging
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

# Try to import markitdown Python API
try:
    from markitdown import MarkItDown
    _HAS_MARKITDOWN = True
except Exception as e:
    logger.warning(f"markitdown Python API not available: {e}")
    _HAS_MARKITDOWN = False

# Try to locate CLI
def _find_markitdown_cli():
    import shutil
    path = shutil.which("markitdown")
    if path:
        return path
    # Common Windows Python Scripts path
    import sys
    scripts = Path(sys.executable).parent / "Scripts" / "markitdown.exe"
    if scripts.exists():
        return str(scripts)
    return None

_MARKITDOWN_CLI = _find_markitdown_cli()


class MarkItDownSkill:
    """Wrapper for markitdown. Falls back to basic extraction if unavailable."""

    def convert(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        logger.info(f"Converting {file_path} to markdown...")

        # Try Python API first
        if _HAS_MARKITDOWN:
            try:
                md = MarkItDown()
                result = md.convert(str(path))
                return result.text_content
            except Exception as e:
                logger.warning(f"markitdown Python API failed: {e}")

        # Try CLI
        if _MARKITDOWN_CLI:
            try:
                result = subprocess.run(
                    [_MARKITDOWN_CLI, str(path)],
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                )
                if result.returncode == 0:
                    return result.stdout
                logger.warning(f"markitdown CLI failed: {result.stderr}")
            except Exception as e:
                logger.warning(f"markitdown CLI error: {e}")

        # Fallback: basic extraction
        return self._fallback_convert(path)

    def _fallback_convert(self, path: Path) -> str:
        """Basic fallback for when markitdown is unavailable."""
        suffix = path.suffix.lower()

        if suffix == ".txt":
            return path.read_text(encoding="utf-8")

        if suffix in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                    parts.append(f"## Slide {i+1}\n\n" + "\n".join(texts))
                return "\n\n".join(parts)
            except Exception as e:
                logger.error(f"Fallback pptx extraction failed: {e}")
                return f"[Error extracting {path.name}: {e}]"

        if suffix in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(str(path))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n\n".join(paragraphs)
            except Exception as e:
                logger.error(f"Fallback docx extraction failed: {e}")
                return f"[Error extracting {path.name}: {e}]"

        if suffix == ".pdf":
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(str(path))
                texts = []
                for page in reader.pages:
                    texts.append(page.extract_text() or "")
                return "\n\n".join(texts)
            except Exception as e:
                logger.error(f"Fallback pdf extraction failed: {e}")
                return f"[Error extracting {path.name}: {e}]"

        return f"[Unsupported file type for fallback extraction: {path.name}]"

    def convert_to_file(self, file_path: str, output_path: str) -> str:
        text = self.convert(file_path)
        Path(output_path).write_text(text, encoding="utf-8")
        return text

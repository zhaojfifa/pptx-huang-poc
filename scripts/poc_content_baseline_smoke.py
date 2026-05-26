#!/usr/bin/env python3
"""
POC 4.5 content-baseline smoke check (P0) — READ-ONLY.

Guards the job_28 content-quality baseline before any agenda / table / custom fix. Given a generated
12-page T5 .pptx, it asserts that body content has NOT collapsed and reports (never cleans) any
template-residue terms.

Read-only: opens the .pptx with python-pptx and inspects text/tables only. It NEVER writes the file,
NEVER cleans residue, NEVER calls the LLM/DB, and makes NO visual-quality claims — it only checks that
text volume / table-cell counts stay at or above conservative baseline thresholds.

Thresholds are derived conservatively from the accepted baselines (cloud job_28 ≈ local job_121:
P4 ≈ 367 text chars, P8 ≈ 30 non-empty table cells, P9 ≈ 24). Defaults sit well below those so a
healthy deck passes and only a real collapse fails. Override via flags if a different template is used.

Usage:
    python -m scripts.poc_content_baseline_smoke [PPTX]            # or a path; default = newest output/job_*_final.pptx
    python -m scripts.poc_content_baseline_smoke --pptx output/job_121_final.pptx
    python -m scripts.poc_content_baseline_smoke --expected-slides 12 --p4-min-chars 150 \
        --p8-min-cells 12 --p9-min-cells 12

Exit code: 0 if PASS or WARN (residue only); 1 if any FAIL (so CI can gate). Never modifies anything.
"""

import argparse
import glob
import os
import sys

# Known T5/T6 template-residue tokens. Reported as WARNINGS only — this script never cleans them
# (global residue cleaning caused the Q3d body-content regression; cleaning is out of scope here).
RESIDUE_TERMS = [
    "穿透式监管", "穿透监管", "VPN", "SD-WAN", "是否模型", "模型场景", "分析场景",
    "全级次", "全链条", "全要素", "国资委", "中国移动", "中移", "CMCC", "OneCity",
]

DEFAULTS = {
    "expected_slides": 12,
    "p4_min_chars": 150,   # baseline ~367; conservative floor
    "p8_min_cells": 12,    # baseline ~30
    "p9_min_cells": 12,    # baseline ~24
}


def _newest_default_pptx() -> str | None:
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    cands = sorted(glob.glob(os.path.join(here, "output", "job_*_final.pptx")),
                   key=os.path.getmtime, reverse=True)
    return cands[0] if cands else None


def _page_metrics(slide):
    """Return (text_chars, text_shapes, table_cells, all_text) for one slide. Read-only."""
    chars = shapes = cells = 0
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                chars += len(t)
                shapes += 1
                texts.append(t)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    ct = cell.text.strip()
                    if ct:
                        cells += 1
                        texts.append(ct)
    return chars, shapes, cells, texts


def run(pptx_path: str, cfg: dict) -> int:
    try:
        from pptx import Presentation
    except ImportError:
        print("FAIL: python-pptx not installed.")
        return 1

    checks = []  # (level, name, detail)  level in {PASS, WARN, FAIL}

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"FAIL: cannot open PPTX: {e}")
        return 1
    checks.append(("PASS", "opens", f"{pptx_path}"))

    slides = list(prs.slides)
    n = len(slides)
    checks.append((
        "PASS" if n == cfg["expected_slides"] else "FAIL",
        "slide_count", f"{n} (expected {cfg['expected_slides']})",
    ))

    # Per-page body metrics (1-based P4/P8/P9 -> 0-based index)
    metrics = {}
    for p in (4, 8, 9):
        if p - 1 < n:
            chars, tshapes, tcells, _ = _page_metrics(slides[p - 1])
            metrics[p] = (chars, tshapes, tcells)
        else:
            metrics[p] = (0, 0, 0)

    p4_chars = metrics[4][0]
    checks.append(("PASS" if p4_chars >= cfg["p4_min_chars"] else "FAIL",
                   "P4_text_chars", f"{p4_chars} (min {cfg['p4_min_chars']}, {metrics[4][1]} text shapes)"))
    for p in (8, 9):
        cells = metrics[p][2]
        key = f"p{p}_min_cells"
        checks.append(("PASS" if cells >= cfg[key] else "FAIL",
                       f"P{p}_table_cells", f"{cells} (min {cfg[key]})"))

    # Contamination — REPORT ONLY (warning), never clean, never fail on it.
    all_text = "\n".join(t for s in slides for _, _, _, txts in [_page_metrics(s)] for t in txts)
    hits = {term: all_text.count(term) for term in RESIDUE_TERMS if term in all_text}
    if hits:
        checks.append(("WARN", "contamination",
                       "; ".join(f"{k}×{v}" for k, v in sorted(hits.items(), key=lambda kv: -kv[1]))))
    else:
        checks.append(("PASS", "contamination", "no residue terms found"))

    # ---- summary ----
    print(f"\nPOC content-baseline smoke · {os.path.basename(pptx_path)}")
    print(f"thresholds: slides=={cfg['expected_slides']} P4>={cfg['p4_min_chars']}ch "
          f"P8>={cfg['p8_min_cells']}cells P9>={cfg['p9_min_cells']}cells")
    print("-" * 64)
    for level, name, detail in checks:
        print(f"  [{level:4}] {name:18} {detail}")
    print("-" * 64)
    n_fail = sum(1 for c in checks if c[0] == "FAIL")
    n_warn = sum(1 for c in checks if c[0] == "WARN")
    overall = "FAIL" if n_fail else ("WARN" if n_warn else "PASS")
    print(f"OVERALL: {overall}  (fail={n_fail} warn={n_warn})")
    if overall == "WARN":
        print("note: WARN = template residue present (advisory). Body thresholds OK; not auto-cleaned.")
    return 1 if n_fail else 0


def main():
    ap = argparse.ArgumentParser(description="Read-only POC content-baseline smoke check (P0).")
    ap.add_argument("pptx", nargs="?", help="Path to a generated .pptx (default: newest output/job_*_final.pptx)")
    ap.add_argument("--pptx", dest="pptx_opt", help="Path to a generated .pptx (alternative to positional)")
    ap.add_argument("--expected-slides", type=int, default=DEFAULTS["expected_slides"])
    ap.add_argument("--p4-min-chars", type=int, default=DEFAULTS["p4_min_chars"])
    ap.add_argument("--p8-min-cells", type=int, default=DEFAULTS["p8_min_cells"])
    ap.add_argument("--p9-min-cells", type=int, default=DEFAULTS["p9_min_cells"])
    args = ap.parse_args()

    pptx_path = args.pptx_opt or args.pptx or _newest_default_pptx()
    if not pptx_path:
        print("FAIL: no PPTX given and no output/job_*_final.pptx found.")
        sys.exit(1)
    if not os.path.exists(pptx_path):
        print(f"FAIL: PPTX not found: {pptx_path}")
        sys.exit(1)

    cfg = {
        "expected_slides": args.expected_slides,
        "p4_min_chars": args.p4_min_chars,
        "p8_min_cells": args.p8_min_cells,
        "p9_min_cells": args.p9_min_cells,
    }
    sys.exit(run(pptx_path, cfg))


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
POC outline→PPT fidelity gate (P1.6) — READ-ONLY.

Given a generated deck and its confirmed_outline.json, reports whether the rendered slide 1 cover and
slide 2 agenda match the confirmed outline, plus the P0 body-content metrics. It NEVER modifies the
deck, NEVER cleans residue, NEVER calls the LLM/DB, and applies NO auto-fix — it only reports and
exits non-zero on a hard body/structure failure so it can gate future changes.

Checks:
  opens, slide_count == expected
  cover_title_match    — visible cover title shape text == confirmed_outline page-1 title
  cover_subtitle_match — cover subtitle shape (名含 副标题/subtitle) text == page-1 first point (skip if no such shape)
  agenda_items_match   — agenda slide contains the deduped body section_titles (PASS=all; WARN=leading subset/overflow)
  P4_text_chars, P8_table_cells, P9_table_cells >= baseline floors (HARD)

Hard gates (exit 1): opens, slide_count, P4/P8/P9 floors. cover/agenda matches are reported; a cover
title MISMATCH is also a hard FAIL (the fill should guarantee it); subtitle/agenda nuances are WARN.

Usage:
    python -m scripts.poc_outline_fidelity_gate [PPTX] [--confirmed-outline PATH]
        [--expected-slides 12] [--p4-min-chars 150] [--p8-min-cells 12] [--p9-min-cells 12]
"""

import argparse
import glob
import json
import os
import re
import sys

DEFAULTS = {"expected_slides": 12, "p4_min_chars": 150, "p8_min_cells": 12, "p9_min_cells": 12}


def _newest_pptx():
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    c = sorted(glob.glob(os.path.join(here, "output", "job_*_final.pptx")),
               key=os.path.getmtime, reverse=True)
    return c[0] if c else None


def _confirmed_for(pptx_path):
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    m = re.search(r"job_(\d+)_final\.pptx$", os.path.basename(pptx_path))
    if not m:
        return None
    p = os.path.join(here, "logs", f"job_{m.group(1)}", "confirmed_outline.json")
    return p if os.path.exists(p) else None


def _is_title_shape(sh):
    name = sh.name or ""
    return ("标题" in name) or ("title" in name.lower())


def _shape_texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh.text_frame.text.strip())
        if sh.has_table:
            for r in sh.table.rows:
                for c in r.cells:
                    if c.text.strip():
                        out.append(c.text.strip())
    return out


def _page_table_cells(slide):
    n = 0
    for sh in slide.shapes:
        if sh.has_table:
            for r in sh.table.rows:
                for c in r.cells:
                    if c.text.strip():
                        n += 1
    return n


def _page_text_chars(slide):
    return sum(len(sh.text_frame.text.strip()) for sh in slide.shapes
               if sh.has_text_frame and sh.text_frame.text.strip())


def run(pptx_path, confirmed_path, cfg):
    from pptx import Presentation
    checks = []  # (level, name, detail)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"FAIL: cannot open PPTX: {e}")
        return 1
    checks.append(("PASS", "opens", os.path.basename(pptx_path)))
    slides = list(prs.slides)
    n = len(slides)
    checks.append(("PASS" if n == cfg["expected_slides"] else "FAIL",
                   "slide_count", f"{n} (expected {cfg['expected_slides']})"))

    outline = {}
    if confirmed_path and os.path.exists(confirmed_path):
        try:
            outline = json.load(open(confirmed_path, encoding="utf-8"))
        except Exception as e:
            checks.append(("WARN", "confirmed_outline", f"unreadable: {e}"))
    else:
        checks.append(("WARN", "confirmed_outline", "not found — fidelity checks skipped"))

    osl = outline.get("slides", []) if outline else []

    # ---- cover fidelity ----
    if osl and n >= 1:
        cov = osl[0]
        want_title = (cov.get("title") or "").strip()
        want_sub = ""
        pts = cov.get("key_points") or []
        if pts:
            want_sub = (pts[0] if isinstance(pts[0], str) else str(pts[0])).strip()
        cover = slides[0]
        tshapes = [sh for sh in cover.shapes if sh.has_text_frame]
        title_sh = next((sh for sh in tshapes if _is_title_shape(sh)), None)
        if title_sh is None and tshapes:
            title_sh = sorted(tshapes, key=lambda s: (s.top or 0))[0]
        got_title = (title_sh.text_frame.text.strip() if title_sh else "")
        if want_title:
            checks.append(("PASS" if got_title == want_title else "FAIL",
                           "cover_title_match", f"want={want_title!r} got={got_title[:40]!r}"))
        sub_sh = next((sh for sh in tshapes if sh is not title_sh
                       and (("副标题" in (sh.name or "")) or ("subtitle" in (sh.name or "").lower()))), None)
        if want_sub and sub_sh is not None:
            got_sub = sub_sh.text_frame.text.strip()
            checks.append(("PASS" if got_sub == want_sub else "WARN",
                           "cover_subtitle_match", f"want={want_sub!r} got={got_sub[:40]!r}"))
        else:
            checks.append(("PASS", "cover_subtitle_match", "no subtitle shape / no cover point (skipped)"))

    # ---- agenda fidelity ----
    if osl and n >= 2:
        body_sections = []
        for c in osl[2:len(osl) - 1]:
            st = (c.get("section_title") or "").strip()
            if st and st not in body_sections:
                body_sections.append(st)
        agenda_blob = "\n".join(_shape_texts(slides[1]))
        present = [s for s in body_sections if s in agenda_blob]
        if not body_sections:
            checks.append(("WARN", "agenda_items_match", "no body section_titles in outline"))
        elif len(present) == len(body_sections):
            checks.append(("PASS", "agenda_items_match", f"{len(present)}/{len(body_sections)} present"))
        elif present and present == body_sections[:len(present)]:
            checks.append(("WARN", "agenda_items_match",
                           f"{len(present)}/{len(body_sections)} present (leading subset — template slot overflow)"))
        elif present:
            checks.append(("WARN", "agenda_items_match", f"{len(present)}/{len(body_sections)} present (out of order)"))
        else:
            checks.append(("FAIL", "agenda_items_match", "0 body section_titles found on agenda slide"))

    # ---- P0 body metrics (HARD) ----
    def metric(p):
        return slides[p - 1] if p - 1 < n else None
    p4 = metric(4)
    p4c = _page_text_chars(p4) if p4 else 0
    checks.append(("PASS" if p4c >= cfg["p4_min_chars"] else "FAIL", "P4_text_chars",
                   f"{p4c} (min {cfg['p4_min_chars']})"))
    for p in (8, 9):
        sl = metric(p)
        cells = _page_table_cells(sl) if sl else 0
        checks.append(("PASS" if cells >= cfg[f"p{p}_min_cells"] else "FAIL",
                       f"P{p}_table_cells", f"{cells} (min {cfg[f'p{p}_min_cells']})"))

    # ---- summary ----
    print(f"\nPOC outline→PPT fidelity gate · {os.path.basename(pptx_path)}")
    print("-" * 64)
    for level, name, detail in checks:
        print(f"  [{level:4}] {name:20} {detail}")
    print("-" * 64)
    n_fail = sum(1 for c in checks if c[0] == "FAIL")
    n_warn = sum(1 for c in checks if c[0] == "WARN")
    overall = "FAIL" if n_fail else ("WARN" if n_warn else "PASS")
    print(f"OVERALL: {overall}  (fail={n_fail} warn={n_warn})")
    return 1 if n_fail else 0


def main():
    ap = argparse.ArgumentParser(description="Read-only outline→PPT fidelity gate (P1.6).")
    ap.add_argument("pptx", nargs="?", help="generated .pptx (default: newest output/job_*_final.pptx)")
    ap.add_argument("--confirmed-outline", help="path to confirmed_outline.json (default: logs/job_<N>/)")
    ap.add_argument("--expected-slides", type=int, default=DEFAULTS["expected_slides"])
    ap.add_argument("--p4-min-chars", type=int, default=DEFAULTS["p4_min_chars"])
    ap.add_argument("--p8-min-cells", type=int, default=DEFAULTS["p8_min_cells"])
    ap.add_argument("--p9-min-cells", type=int, default=DEFAULTS["p9_min_cells"])
    args = ap.parse_args()

    pptx_path = args.pptx or _newest_pptx()
    if not pptx_path or not os.path.exists(pptx_path):
        print(f"FAIL: PPTX not found: {pptx_path}")
        sys.exit(1)
    confirmed = args.confirmed_outline or _confirmed_for(pptx_path)
    cfg = {"expected_slides": args.expected_slides, "p4_min_chars": args.p4_min_chars,
           "p8_min_cells": args.p8_min_cells, "p9_min_cells": args.p9_min_cells}
    sys.exit(run(pptx_path, confirmed, cfg))


if __name__ == "__main__":
    main()
